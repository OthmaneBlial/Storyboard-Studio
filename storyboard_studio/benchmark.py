"""Reproducible synthetic benchmark for decision-story release candidates."""

from __future__ import annotations

import json
import re
import shutil
from collections import defaultdict
from collections.abc import Mapping
from pathlib import Path
from typing import Any, Literal

from pptx import Presentation
from pydantic import BaseModel, ConfigDict, Field, field_validator, model_validator

from ai_helper import GenerationRun, generate_ppt_content_run
from generate_pptx import create_presentation
from schemas import DecisionBriefV2, PresentationPayload, StoryDocumentV2
from storyboard_studio import __version__
from storyboard_studio.doctor import diagnose_story
from storyboard_studio.evidence import evidence_coverage
from storyboard_studio.layout import analyze_overflow, load_layout_contract
from storyboard_studio.providers import EXCLUDED_FIELDS
from storyboard_studio.receipt import digest_file, digest_value
from storyboard_studio.semantic import block_plain_text, normalize_content_block
from storyboard_studio.story import build_decision_story, migrate_presentation_v1

WORD_RE = re.compile(r"[a-z0-9]+", re.IGNORECASE)
NUMBER_RE = re.compile(r"(?<![\w/])(?:[$€£]\s*)?\d+(?:[.,]\d+)?%?(?![\w/])")
CATEGORIES = (
    "content",
    "design",
    "coherence",
    "editability",
    "provenance",
    "privacy",
    "reproducibility",
)
KNOWN_LIMITATIONS = (
    "The deterministic rubric checks published structural signals; it is not a human or "
    "vision-language-model judgment of presentation quality.",
    "The baseline inspects PPTX structure but does not run PowerPoint, Keynote, or Google Slides; "
    "viewer compatibility remains a separate release gate.",
    "All briefs and evidence are synthetic fixtures, so this benchmark is not user research or "
    "proof of real-world usefulness.",
    "Unless provider network access is explicitly enabled, the optional-provider lane records a "
    "not-configured fallback and does not measure model quality.",
    "Reproducibility compares canonical story semantics; PPTX ZIP bytes may differ because package "
    "metadata and compression are not a semantic contract.",
)


class BenchmarkModel(BaseModel):
    model_config = ConfigDict(extra="forbid", str_strip_whitespace=True)


class BenchmarkCase(BenchmarkModel):
    id: str = Field(pattern=r"^[a-z0-9][a-z0-9-]{2,63}$")
    theme: Literal["midnight", "glacier", "ember", "forest", "royal", "sakura"]
    copy_density_risks: list[Literal["current_context", "constraints", "option_descriptions", "trade_offs"]]
    expected_author_checked_claims: int = Field(ge=0)
    expected_unresolved_claims: int = Field(ge=0)
    brief: DecisionBriefV2


class BenchmarkSuite(BenchmarkModel):
    schema_version: Literal["1"] = "1"
    id: str = Field(pattern=r"^[a-z0-9][a-z0-9-]{2,63}$")
    title: str = Field(min_length=3, max_length=120)
    license: Literal["CC0-1.0"]
    synthetic_only: Literal[True]
    viewer_constraints: list[str] = Field(min_length=5)
    expected_story_roles: list[str] = Field(min_length=5, max_length=5)
    expected_blocks: list[str] = Field(min_length=5, max_length=5)
    cases: list[BenchmarkCase] = Field(min_length=1, max_length=50)

    @field_validator("cases")
    @classmethod
    def unique_case_ids(cls, cases: list[BenchmarkCase]) -> list[BenchmarkCase]:
        identifiers = [case.id for case in cases]
        if len(identifiers) != len(set(identifiers)):
            raise ValueError("Benchmark case ids must be unique.")
        return cases

    @model_validator(mode="after")
    def expected_claim_counts_are_possible(self) -> BenchmarkSuite:
        for case in self.cases:
            story = build_decision_story(case.brief, case.theme)
            summary = evidence_coverage(story.presentation)["summary"]
            expected = (case.expected_author_checked_claims, case.expected_unresolved_claims)
            actual = (summary["author_checked_claims"], summary["unresolved_claims"])
            if expected != actual:
                raise ValueError(
                    f"Case {case.id!r} expects evidence counts {expected}, but produces {actual}."
                )
        return self


def _reject_duplicate_keys(pairs: list[tuple[str, Any]]) -> dict[str, Any]:
    result: dict[str, Any] = {}
    for key, value in pairs:
        if key in result:
            raise ValueError(f"Duplicate JSON key {key!r} is not allowed.")
        result[key] = value
    return result


def load_benchmark_suite(path: str | Path) -> BenchmarkSuite:
    source = Path(path).expanduser().resolve()
    try:
        value = json.loads(source.read_text(encoding="utf-8"), object_pairs_hook=_reject_duplicate_keys)
    except (OSError, json.JSONDecodeError, ValueError) as exc:
        raise ValueError(f"Invalid benchmark suite {source}: {exc}") from exc
    return BenchmarkSuite.model_validate(value)


def _write_json(path: Path, value: object) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(
        json.dumps(value, indent=2, ensure_ascii=False, default=str) + "\n",
        encoding="utf-8",
    )


def _words(value: str) -> set[str]:
    return {word.lower() for word in WORD_RE.findall(value) if len(word) > 2}


def _story_text(story: StoryDocumentV2) -> str:
    parts = [story.presentation.title, story.presentation.subtitle]
    for slide in story.presentation.slides:
        parts.extend((slide.title, slide.content))
        parts.append(block_plain_text(normalize_content_block(slide.model_dump(mode="json"))))
    return " ".join(parts)


def _term_coverage(expected: str, actual: str) -> float:
    expected_words = _words(expected)
    return len(expected_words & _words(actual)) / len(expected_words) if expected_words else 1.0


def _criterion(name: str, points: float, maximum: float, evidence: str) -> dict[str, object]:
    return {
        "name": name,
        "points": round(max(0.0, min(points, maximum)), 2),
        "maximum": maximum,
        "evidence": evidence,
    }


def _pptx_inspection(path: Path) -> dict[str, object]:
    deck = Presentation(path)
    text_shapes = pictures = full_slide_pictures = 0
    for slide in deck.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                text_shapes += 1
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                pictures += 1
                if shape.width >= deck.slide_width * 0.95 and shape.height >= deck.slide_height * 0.95:
                    full_slide_pictures += 1
    ratio = deck.slide_width / deck.slide_height
    return {
        "slides_including_title": len(deck.slides),
        "width_emu": deck.slide_width,
        "height_emu": deck.slide_height,
        "aspect_ratio": round(ratio, 6),
        "is_16_9": abs(ratio - 16 / 9) <= 0.002,
        "text_shapes": text_shapes,
        "pictures": pictures,
        "full_slide_pictures": full_slide_pictures,
    }


def _provider_story(
    case: BenchmarkCase, provider: str, environment: dict[str, str] | None
) -> tuple[StoryDocumentV2, dict[str, object]]:
    brief = case.brief
    supporting_brief = " | ".join(
        (
            f"Audience: {brief.audience}",
            f"Outcome: {brief.desired_outcome}",
            f"Constraints: {'; '.join(brief.constraints)}",
            f"Options: {'; '.join(option.title for option in brief.options)}",
            f"Trade-offs: {'; '.join(brief.trade_offs)}",
            f"Owner: {brief.owner}",
            f"Next: {brief.next_step}",
        )
    )
    configs = [
        {"focus": role, "block": block, "layout": "focus" if index == 2 else "right"}
        for index, (role, block) in enumerate(
            zip(
                ("decision context", "constraints", "options", "trade-off", "owned next step"),
                ("standard", "standard", "comparison", "decision", "timeline"),
                strict=True,
            )
        )
    ]
    run: GenerationRun = generate_ppt_content_run(
        brief.decision,
        5,
        supporting_brief,
        configs,
        use_ai=True,
        provider=provider,
        environment=environment,
    )
    story = migrate_presentation_v1(PresentationPayload.model_validate(run.presentation))
    story = story.model_copy(update={"planner": run.source, "provider_warning": run.warning or ""})
    return story, run.provider


def _local_provider_metadata() -> dict[str, object]:
    return {
        "selected": "local",
        "used": "local",
        "model": "deterministic-v1",
        "used_model": "deterministic-v1",
        "network_boundary": "offline",
        "network_status": "offline",
        "transferred_fields": [],
        "excluded_fields": list(EXCLUDED_FIELDS),
        "fallback_reason": None,
    }


def _score_story(
    case: BenchmarkCase,
    suite: BenchmarkSuite,
    story: StoryDocumentV2,
    provider: Mapping[str, object],
    pptx: Path,
    repeat_digest: str | None,
) -> dict[str, object]:
    doctor = diagnose_story(story)
    coverage = evidence_coverage(story.presentation)
    layout = analyze_overflow(story.presentation.model_dump(mode="json"), load_layout_contract())
    inspection = _pptx_inspection(pptx)
    text = _story_text(story)
    finding_codes = {finding["code"] for finding in doctor["findings"]}
    blocks = [slide.block for slide in story.presentation.slides]
    brief = case.brief

    decision_coverage = _term_coverage(brief.decision, text)
    audience_coverage = max(_term_coverage(brief.audience, text), _term_coverage(brief.desired_outcome, text))
    option_matches = sum(option.title.lower() in text.lower() for option in brief.options)
    owner_coverage = max(_term_coverage(brief.owner, text), _term_coverage(brief.next_step, text))
    expected_numbers = set(
        NUMBER_RE.findall(json.dumps(brief.model_dump(mode="json"), ensure_ascii=False, default=str))
    )
    expected_numbers.update({"1", "2", "3", "01", "02", "03"})
    unexpected_numbers = set(NUMBER_RE.findall(text)) - expected_numbers
    categories = {
        "content": [
            _criterion(
                "decision term coverage",
                6 if decision_coverage >= 0.6 else 3 if decision_coverage >= 0.3 else 0,
                6,
                f"{decision_coverage:.0%} of decision terms appear in the story.",
            ),
            _criterion(
                "audience or outcome coverage",
                4 if audience_coverage >= 0.5 else 2 if audience_coverage >= 0.25 else 0,
                4,
                f"Best coverage is {audience_coverage:.0%}.",
            ),
            _criterion(
                "option title preservation",
                4 * option_matches / len(brief.options),
                4,
                f"{option_matches}/{len(brief.options)} option titles appear verbatim.",
            ),
            _criterion(
                "owner or next-step coverage",
                4 if owner_coverage >= 0.5 else 2 if owner_coverage >= 0.25 else 0,
                4,
                f"Best coverage is {owner_coverage:.0%}.",
            ),
            _criterion(
                "numeric restraint",
                2 if not unexpected_numbers else 0,
                2,
                "No numeric tokens were introduced beyond the brief or structural enumerators."
                if not unexpected_numbers
                else f"Unexpected numeric tokens: {sorted(unexpected_numbers)}.",
            ),
        ],
        "design": [
            _criterion(
                "shared overflow budget",
                8 if not layout["findings"] else max(0, 8 - 2 * len(layout["findings"])),
                8,
                f"{len(layout['findings'])} deterministic overflow finding(s).",
            ),
            _criterion(
                "semantic block variety",
                4 if len(set(blocks)) >= 4 else len(set(blocks)),
                4,
                f"{len(set(blocks))} distinct block role(s): {blocks}.",
            ),
            _criterion(
                "native editable structure",
                5 if inspection["text_shapes"] and not inspection["full_slide_pictures"] else 0,
                5,
                f"{inspection['text_shapes']} text shapes; "
                f"{inspection['full_slide_pictures']} full-slide pictures.",
            ),
            _criterion(
                "16:9 canvas",
                3 if inspection["is_16_9"] else 0,
                3,
                f"Measured aspect ratio {inspection['aspect_ratio']}.",
            ),
        ],
        "coherence": [
            _criterion(
                "expected story-role sequence",
                8 if blocks == suite.expected_blocks else 4 if len(blocks) == 5 else 0,
                8,
                f"Observed blocks: {blocks}.",
            ),
            _criterion(
                "explicit decision",
                4 if "decision.missing" not in finding_codes else 0,
                4,
                "Doctor decision.missing is absent."
                if "decision.missing" not in finding_codes
                else "Doctor reported decision.missing.",
            ),
            _criterion(
                "visible trade-off",
                4 if "tradeoff.missing" not in finding_codes else 0,
                4,
                "Doctor tradeoff.missing is absent."
                if "tradeoff.missing" not in finding_codes
                else "Doctor reported tradeoff.missing.",
            ),
            _criterion(
                "owned next action",
                4 if "action.missing" not in finding_codes else 0,
                4,
                "Doctor action.missing is absent."
                if "action.missing" not in finding_codes
                else "Doctor reported action.missing.",
            ),
        ],
        "editability": [
            _criterion(
                "story text exported as text frames",
                6 if inspection["text_shapes"] >= len(story.presentation.slides) else 0,
                6,
                f"{inspection['text_shapes']} text shapes for "
                f"{len(story.presentation.slides)} content slides.",
            ),
            _criterion(
                "no rasterized full-slide content",
                4 if not inspection["full_slide_pictures"] else 0,
                4,
                f"{inspection['full_slide_pictures']} full-slide picture(s).",
            ),
        ],
    }
    evidence_summary = coverage["summary"]
    evidence_exact = (
        evidence_summary["author_checked_claims"] == case.expected_author_checked_claims
        and evidence_summary["unresolved_claims"] == case.expected_unresolved_claims
    )
    categories["provenance"] = [
        _criterion(
            "claim-level coverage report",
            3 if evidence_summary["claims"] else 0,
            3,
            f"{evidence_summary['claims']} claims were enumerated.",
        ),
        _criterion(
            "unresolved gaps remain explicit",
            3 if evidence_summary["unresolved_claims"] else 0,
            3,
            f"{evidence_summary['unresolved_claims']} claims remain unresolved.",
        ),
        _criterion(
            "expected evidence state preserved",
            4 if evidence_exact else 0,
            4,
            f"Expected/actual checked {case.expected_author_checked_claims}/"
            f"{evidence_summary['author_checked_claims']}; unresolved "
            f"{case.expected_unresolved_claims}/{evidence_summary['unresolved_claims']}.",
        ),
    ]
    network_status = str(provider.get("network_status", "unknown"))
    boundary = str(provider.get("network_boundary", "unknown"))
    network_matches = (
        (boundary == "offline" and network_status == "offline")
        or (boundary == "loopback-only" and network_status.startswith(("loopback", "not-sent")))
        or (boundary == "external-provider" and network_status.startswith(("external", "not-sent")))
    )
    excluded = set(provider.get("excluded_fields", []))
    categories["privacy"] = [
        _criterion(
            "declared network boundary",
            4 if network_matches else 0,
            4,
            f"Boundary {boundary}; observed status {network_status}.",
        ),
        _criterion(
            "private payload fields excluded",
            4 if set(EXCLUDED_FIELDS) <= excluded else 0,
            4,
            f"Excluded fields: {sorted(excluded)}.",
        ),
        _criterion(
            "no remote presentation assets",
            2 if not story.presentation.assets else 0,
            2,
            f"{len(story.presentation.assets)} local asset record(s); remote assets are unsupported.",
        ),
    ]
    story_digest = digest_value(story.model_dump(mode="json"))
    categories["reproducibility"] = [
        _criterion(
            "canonical story repeatability",
            6 if repeat_digest == story_digest else 0,
            6,
            "A second generation produced the same canonical story digest."
            if repeat_digest == story_digest
            else "Repeatability was skipped for a live provider or the semantic digest changed.",
        ),
        _criterion(
            "tool and schema versions recorded",
            4,
            4,
            f"Storyboard Studio {__version__}; story schema {story.schema_version}.",
        ),
    ]
    category_scores = {
        name: round(sum(float(item["points"]) for item in criteria), 2)
        for name, criteria in categories.items()
    }
    return {
        "schema_version": "1",
        "case_id": case.id,
        "mode": "local" if provider.get("selected") == "local" else "optional-provider",
        "score": round(sum(category_scores.values()), 2),
        "maximum": 100,
        "category_scores": category_scores,
        "criteria": categories,
        "provider": dict(provider),
        "story_sha256": story_digest,
        "pptx_sha256": digest_file(pptx),
        "pptx_inspection": inspection,
        "copy_density_risks": case.copy_density_risks,
        "viewer_constraints": suite.viewer_constraints,
    }


def _repeat_story(
    case: BenchmarkCase,
    mode: str,
    provider_id: str,
    provider_metadata: Mapping[str, object],
    environment: dict[str, str] | None,
) -> str | None:
    if mode == "local":
        repeated = build_decision_story(case.brief, case.theme)
        return digest_value(repeated.model_dump(mode="json"))
    if provider_metadata.get("used") != "local":
        return None
    repeated, _ = _provider_story(case, provider_id, environment)
    return digest_value(repeated.model_dump(mode="json"))


def _case_mode_artifacts(
    case: BenchmarkCase,
    suite: BenchmarkSuite,
    mode: Literal["local", "optional-provider"],
    output_dir: Path,
    provider_id: str,
    environment: dict[str, str] | None,
) -> dict[str, object]:
    if mode == "local":
        story, provider = build_decision_story(case.brief, case.theme), _local_provider_metadata()
    else:
        story, provider = _provider_story(case, provider_id, environment)
    destination = output_dir / "cases" / case.id / mode
    destination.mkdir(parents=True, exist_ok=True)
    pptx_path = destination / "deck.pptx"
    _write_json(destination / "story.json", story.model_dump(mode="json"))
    create_presentation(story.presentation.model_dump(mode="json"), pptx_path)
    _write_json(destination / "doctor.json", diagnose_story(story))
    _write_json(destination / "evidence.json", evidence_coverage(story.presentation))
    _write_json(
        destination / "layout.json",
        analyze_overflow(story.presentation.model_dump(mode="json"), load_layout_contract()),
    )
    _write_json(destination / "provider.json", provider)
    score = _score_story(
        case,
        suite,
        story,
        provider,
        pptx_path,
        _repeat_story(case, mode, provider_id, provider, environment),
    )
    _write_json(destination / "score.json", score)
    return score


def _summarize(results: list[dict[str, object]]) -> dict[str, object]:
    by_mode: dict[str, list[dict[str, object]]] = defaultdict(list)
    for result in results:
        by_mode[str(result["mode"])].append(result)
    modes: dict[str, object] = {}
    for mode, rows in sorted(by_mode.items()):
        modes[mode] = {
            "cases": len(rows),
            "average_score": round(sum(float(row["score"]) for row in rows) / len(rows), 2),
            "category_averages": {
                category: round(sum(float(row["category_scores"][category]) for row in rows) / len(rows), 2)
                for category in CATEGORIES
            },
        }
    return {"runs": len(results), "modes": modes}


def compare_reports(current: Mapping[str, object], baseline: Mapping[str, object]) -> dict[str, object]:
    baseline_rows = {
        (str(row["case_id"]), str(row["mode"])): row
        for row in baseline.get("results", [])
        if isinstance(row, Mapping)
    }
    comparisons: list[dict[str, object]] = []
    regressions: list[dict[str, object]] = []
    for row in current.get("results", []):
        if not isinstance(row, Mapping):
            continue
        key = (str(row["case_id"]), str(row["mode"]))
        previous = baseline_rows.get(key)
        if previous is None:
            comparisons.append({"case_id": key[0], "mode": key[1], "status": "new"})
            continue
        current_provider, previous_provider = row.get("provider", {}), previous.get("provider", {})
        if (
            isinstance(current_provider, Mapping)
            and isinstance(previous_provider, Mapping)
            and (
                current_provider.get("used") != previous_provider.get("used")
                or current_provider.get("used_model") != previous_provider.get("used_model")
            )
        ):
            comparisons.append({"case_id": key[0], "mode": key[1], "status": "provider-not-comparable"})
            continue
        delta = round(float(row["score"]) - float(previous["score"]), 2)
        comparison = {"case_id": key[0], "mode": key[1], "status": "compared", "delta": delta}
        comparisons.append(comparison)
        if delta < 0:
            regressions.append(comparison)
    return {
        "baseline_release": baseline.get("release", "unknown"),
        "comparisons": comparisons,
        "regressions": regressions,
        "status": "regressed" if regressions else "no-regression",
    }


def _markdown_report(report: Mapping[str, object]) -> str:
    lines = [
        f"# Storyboard benchmark — {report['release']}",
        "",
        f"Suite: `{report['suite_id']}` · Storyboard Studio `{report['tool_version']}` · "
        f"{report['summary']['runs']} raw runs",
        "",
        "| Mode | Average | Content | Design | Coherence | Editability | Provenance | Privacy | "
        "Reproducibility |",
        "| --- | ---: | ---: | ---: | ---: | ---: | ---: | ---: | ---: |",
    ]
    for mode, summary in report["summary"]["modes"].items():
        scores = summary["category_averages"]
        lines.append(
            f"| {mode} | {summary['average_score']}/100 | {scores['content']}/20 | "
            f"{scores['design']}/20 | {scores['coherence']}/20 | "
            f"{scores['editability']}/10 | {scores['provenance']}/10 | "
            f"{scores['privacy']}/10 | {scores['reproducibility']}/10 |"
        )
    lines.extend(["", "## Raw results", ""])
    lines.extend(
        f"- `{row['case_id']}` / `{row['mode']}`: **{row['score']}/100** "
        f"(provider used: `{row['provider']['used']}`, network: "
        f"`{row['provider']['network_status']}`)"
        for row in report["results"]
    )
    lines.extend(["", "## Known limitations", ""])
    lines.extend(f"- {limitation}" for limitation in report["limitations"])
    if report.get("regression"):
        lines.extend(
            [
                "",
                "## Regression comparison",
                "",
                f"Status: **{report['regression']['status']}** against "
                f"`{report['regression']['baseline_release']}`.",
            ]
        )
    lines.extend(
        [
            "",
            "> Scores describe this fixture and rubric only. They do not establish factual truth, "
            "real-user value, or universal visual quality.",
            "",
        ]
    )
    return "\n".join(lines)


def run_benchmark(
    suite_path: str | Path,
    output_dir: str | Path,
    *,
    release: str,
    optional_provider: Literal["gemini", "openai-compatible"] = "openai-compatible",
    allow_provider_network: bool = False,
    baseline_path: str | Path | None = None,
    overwrite: bool = False,
) -> dict[str, object]:
    suite_source = Path(suite_path).expanduser().resolve()
    destination = Path(output_dir).expanduser().resolve()
    suite = load_benchmark_suite(suite_source)
    if destination.exists() and any(destination.iterdir()):
        if not overwrite:
            raise ValueError(f"Benchmark output directory is not empty: {destination}")
        known_outputs = {"cases", "report.json", "report.md", "manifest.json"}
        unrelated = sorted(path.name for path in destination.iterdir() if path.name not in known_outputs)
        if unrelated:
            raise ValueError(
                f"Refusing to overwrite unrelated files in {destination}: {', '.join(unrelated)}"
            )
        for path in destination.iterdir():
            if path.is_dir():
                shutil.rmtree(path)
            else:
                path.unlink()
    destination.mkdir(parents=True, exist_ok=True)
    environment = None if allow_provider_network else {}
    results = []
    for case in suite.cases:
        for mode in ("local", "optional-provider"):
            results.append(
                _case_mode_artifacts(case, suite, mode, destination, optional_provider, environment)
            )
    report: dict[str, object] = {
        "schema_version": "1",
        "release": release,
        "suite_id": suite.id,
        "suite_sha256": digest_file(suite_source),
        "tool_version": __version__,
        "provider_network_allowed": allow_provider_network,
        "optional_provider": optional_provider,
        "summary": _summarize(results),
        "limitations": list(KNOWN_LIMITATIONS),
        "results": results,
    }
    if baseline_path:
        baseline_source = Path(baseline_path).expanduser().resolve()
        try:
            baseline = json.loads(baseline_source.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError) as exc:
            raise ValueError(f"Could not read benchmark baseline {baseline_source}: {exc}") from exc
        report["regression"] = compare_reports(report, baseline)
    _write_json(destination / "report.json", report)
    (destination / "report.md").write_text(_markdown_report(report), encoding="utf-8")
    manifest = {
        str(path.relative_to(destination)): digest_file(path)
        for path in sorted(destination.rglob("*"))
        if path.is_file() and path.name != "manifest.json"
    }
    _write_json(destination / "manifest.json", {"sha256": manifest})
    return report

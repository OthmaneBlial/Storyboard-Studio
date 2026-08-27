import json
from pathlib import Path

from pptx import Presentation

from storyboard_studio.cli import main
from storyboard_studio.resources import web_root


def test_packaged_web_application_is_complete():
    root = web_root()
    assert (root / "index.html").is_file()
    assert (root / "static" / "app.js").is_file()
    assert (root / "static" / "app.css").is_file()
    assert (root / "static" / "favicon.svg").is_file()
    package_root = root.parent
    assert (package_root / "data" / "storyboard-v1.json").is_file()
    assert (package_root / "data" / "story-v2.json").is_file()
    assert (package_root / "data" / "decision-brief.story.json").is_file()
    assert (package_root / "data" / "template-catalog.json").is_file()
    assert (package_root / "data" / "storyboard-tokens.json").is_file()


def test_demo_and_export_commands_create_editable_powerpoints(tmp_path: Path):
    demo = tmp_path / "demo.pptx"
    exported = tmp_path / "exported.pptx"

    assert main(["demo", "--output", str(demo)]) == 0
    assert (
        main(
            [
                "export",
                "--input",
                "examples/product-brief.json",
                "--output",
                str(exported),
            ]
        )
        == 0
    )

    assert len(Presentation(demo).slides) == 6
    assert len(Presentation(exported).slides) == 4


def test_doctor_command_writes_stable_json(tmp_path: Path):
    output = tmp_path / "doctor.json"
    assert (
        main(
            [
                "doctor",
                "examples/product-brief.json",
                "--format",
                "json",
                "--output",
                str(output),
            ]
        )
        == 0
    )

    report = json.loads(output.read_text(encoding="utf-8"))
    assert report["schema_version"] == "1"
    assert report["summary"]["slides"] == 3
    assert "does not verify factual truth" in report["disclaimer"]


def test_compile_command_creates_versioned_decision_story(tmp_path: Path):
    output = tmp_path / "decision.story.json"
    assert (
        main(
            [
                "compile",
                "--input",
                "examples/briefs/onboarding-decision.json",
                "--output",
                str(output),
                "--theme",
                "forest",
            ]
        )
        == 0
    )
    story = json.loads(output.read_text(encoding="utf-8"))
    assert story["schema_version"] == "2"
    assert story["kind"] == "decision-brief"
    assert story["presentation"]["theme"] == "forest"


def test_templates_command_hides_dormant_workflows_by_default(tmp_path: Path):
    launched = tmp_path / "launched.json"
    complete = tmp_path / "complete.json"

    assert main(["templates", "--format", "json", "--output", str(launched)]) == 0
    assert main(["templates", "--all", "--format", "json", "--output", str(complete)]) == 0

    launched_ids = [item["id"] for item in json.loads(launched.read_text())["templates"]]
    complete_items = json.loads(complete.read_text())["templates"]
    assert launched_ids == ["decision-brief"]
    assert [item["id"] for item in complete_items] == [
        "decision-brief",
        "project-alignment",
        "proposal",
        "incident-retrospective",
    ]
    assert all(item["status"] == "dormant" for item in complete_items[1:])


def test_brand_kit_and_layout_preflight_commands_are_local_and_machine_readable(tmp_path: Path):
    kit_result = tmp_path / "brand-kit.json"
    preflight_result = tmp_path / "preflight.json"

    assert main(["brand-kit", "themes/brand-kit.example.json", "--output", str(kit_result)]) == 0
    assert (
        main(
            [
                "preflight",
                "examples/product-brief.json",
                "--output",
                str(preflight_result),
            ]
        )
        == 0
    )

    assert json.loads(kit_result.read_text())["status"] == "valid"
    assert json.loads(preflight_result.read_text())["layout_schema"] == "2"


def test_evidence_report_and_citations_export_are_explicit(tmp_path: Path):
    report_path = tmp_path / "evidence.json"
    deck_path = tmp_path / "evidence.pptx"

    assert (
        main(
            [
                "evidence",
                "examples/fixtures/evidence-edge-cases.json",
                "--output",
                str(report_path),
            ]
        )
        == 0
    )
    assert (
        main(
            [
                "export",
                "--input",
                "examples/fixtures/evidence-edge-cases.json",
                "--output",
                str(deck_path),
                "--citations",
            ]
        )
        == 0
    )

    report = json.loads(report_path.read_text(encoding="utf-8"))
    assert report["summary"]["unresolved_claims"] > 0
    assert len(Presentation(deck_path).slides) == 5

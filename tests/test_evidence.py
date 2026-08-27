import json
from pathlib import Path

import pytest
from pptx import Presentation
from pydantic import ValidationError

from generate_pptx import create_presentation
from outline_markdown import markdown_to_presentation, presentation_to_markdown
from schemas import PresentationPayload, SourceReference
from storyboard_studio.doctor import diagnose_presentation
from storyboard_studio.evidence import approved_citations, evidence_coverage
from storyboard_studio.receipt import create_receipt
from storyboard_studio.story import migrate_presentation_v1


def evidence_fixture() -> PresentationPayload:
    return PresentationPayload.model_validate(
        json.loads(Path("examples/fixtures/evidence-edge-cases.json").read_text(encoding="utf-8"))
    )


def test_claim_coverage_never_treats_a_url_as_verification():
    payload = evidence_fixture()
    report = evidence_coverage(payload)
    unresolved_url_claim = next(
        claim for claim in report["claims"] if claim["slide_number"] == 2 and claim["id"] == "summary"
    )

    assert unresolved_url_claim["has_url"] is True
    assert unresolved_url_claim["status"] == "linked-unresolved"
    assert report["summary"]["author_checked_claims"] < report["summary"]["claims"]
    assert {entry["label"] for entry in approved_citations(payload)} == {
        "Étude d’usage — café",
        "Decision log 27 août",
    }


def test_malicious_urls_long_evidence_and_checked_sources_are_rejected():
    rejected = json.loads(Path("examples/fixtures/evidence-invalid-urls.json").read_text(encoding="utf-8"))[
        "rejected"
    ]
    for url in rejected:
        with pytest.raises(ValidationError):
            SourceReference(label="Unsafe", url=url)

    with pytest.raises(ValidationError, match="300"):
        SourceReference(label="Too long", evidence="x" * 301)
    with pytest.raises(ValidationError, match="owner and checked date"):
        SourceReference(
            label="Incomplete approval",
            url="https://example.org/report",
            review_status="author-checked",
        )
    with pytest.raises(ValidationError, match="relative POSIX"):
        SourceReference(label="Traversal", local_reference="../private/notes.md")


def test_unicode_and_missing_owner_remain_visible_to_the_doctor():
    payload = evidence_fixture()
    report = diagnose_presentation(payload)

    assert "Étude d’usage — café" in payload.slides[0].sources[0].label
    assert any(finding["code"] == "evidence.owner-missing" for finding in report["findings"])
    assert any(finding["code"] == "evidence.url-unchecked" for finding in report["findings"])
    assert report["summary"]["unresolved_claims"] > 0


def test_citations_appendix_is_native_author_checked_and_keeps_notes(tmp_path: Path):
    payload = evidence_fixture()
    exported = Presentation(create_presentation(payload.model_dump(), tmp_path / "evidence.pptx"))
    appendix = exported.slides[-1]
    visible = "\n".join(shape.text for shape in appendix.shapes if shape.has_text_frame)

    assert len(exported.slides) == 5
    assert any(shape.name.startswith("citations.entry") for shape in appendix.shapes)
    assert "Étude d’usage" in visible
    assert "Decision log" in visible
    assert "Unreviewed résumé link" not in visible
    assert "truth not independently verified" in appendix.notes_slide.notes_text_frame.text
    first_notes = exported.slides[1].notes_slide.notes_text_frame.text
    assert "checked 2026-08-27" in first_notes
    assert "claims summary, block-1" in first_notes


def test_sources_survive_markdown_migration_reorder_duplicate_and_receipt(tmp_path: Path):
    payload = evidence_fixture()
    markdown = presentation_to_markdown(payload.model_dump(mode="json"))
    restored = PresentationPayload.model_validate(markdown_to_presentation(markdown))

    assert restored.model_dump(mode="json") == payload.model_dump(mode="json")
    original_sources = payload.slides[0].sources
    reordered = payload.model_copy(deep=True)
    reordered.slides = [reordered.slides[1], reordered.slides[0], reordered.slides[2]]
    for index, slide in enumerate(reordered.slides, start=1):
        slide.slide_number = index
    duplicate = reordered.slides[1].model_copy(deep=True)
    assert duplicate.sources[0].model_dump(mode="json") == original_sources[0].model_dump(mode="json")

    story = migrate_presentation_v1(payload)
    story_path = tmp_path / "evidence.story.json"
    pptx_path = create_presentation(payload.model_dump(), tmp_path / "evidence.pptx")
    story_path.write_text(story.model_dump_json(indent=2), encoding="utf-8")
    receipt = create_receipt(story, story_path, pptx_path)

    assert receipt["source_provenance"][0]["url"].startswith("https://")
    assert receipt["evidence_coverage"]["summary"]["claims"] > 3
    assert len(receipt["approved_citations"]) == 2
    assert receipt["citations_appendix"] is True

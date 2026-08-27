import json
from pathlib import Path

from pptx import Presentation

from ai_helper import build_local_presentation
from generate_pptx import THEMES, create_presentation


def test_renderer_creates_an_editable_title_and_content_slides(tmp_path: Path):
    data = build_local_presentation("A reliable onboarding experience", 3)
    data["theme"] = "forest"
    destination = create_presentation(data, tmp_path / "deck.pptx")

    assert destination.is_file()
    exported = Presentation(destination)
    assert len(exported.slides) == 4
    assert exported.slide_width > exported.slide_height
    assert any(
        "Reliable Onboarding" in shape.text for shape in exported.slides[0].shapes if shape.has_text_frame
    )


def test_all_public_themes_can_render(tmp_path: Path):
    for theme in THEMES:
        data = build_local_presentation("Theme test", 3)
        data["theme"] = theme
        assert create_presentation(data, tmp_path / f"{theme}.pptx").is_file()


def test_edge_fixture_preserves_unicode_and_all_layouts(tmp_path: Path):
    fixture = json.loads(Path("examples/fixtures/edge-cases.json").read_text(encoding="utf-8"))
    destination = create_presentation(fixture, tmp_path / "edge-cases.pptx")
    exported = Presentation(destination)
    texts = "\n".join(
        shape.text for slide in exported.slides for shape in slide.shapes if shape.has_text_frame
    )

    assert len(exported.slides) == 4
    assert "Café déjà vu" in texts
    assert "Unicode text" in texts
    assert "Long copy is a rendering risk" in texts


def test_export_has_stable_editable_core_properties(tmp_path: Path):
    data = build_local_presentation("Core property test", 3)
    exported = Presentation(create_presentation(data, tmp_path / "properties.pptx"))

    assert exported.core_properties.author == "Storyboard Studio"
    assert exported.core_properties.title == "Core Property Test"
    assert all(
        shape.has_text_frame for slide in exported.slides for shape in slide.shapes if shape.has_text_frame
    )


def test_sources_and_speaker_notes_are_native_notes(tmp_path: Path):
    data = build_local_presentation("Evidence-aware brief", 3)
    data["slides"][0]["sources"] = [{"label": "Internal brief", "evidence": "Author notes", "owner": "Team"}]
    data["slides"][0]["speaker_notes"] = "Check this draft before sharing."
    exported = Presentation(create_presentation(data, tmp_path / "notes.pptx"))
    notes = exported.slides[1].notes_slide.notes_text_frame.text

    assert "Check this draft" in notes
    assert "Internal brief" in notes
    assert "author-supplied; not verified" in notes


def test_all_semantic_blocks_render_as_distinct_native_structures_in_dark_and_light(
    tmp_path: Path,
):
    fixture = json.loads(Path("examples/fixtures/semantic-blocks.json").read_text(encoding="utf-8"))
    expected_names = [
        "semantic.standard.point.1",
        "semantic.comparison.side.1",
        "semantic.decision.statement",
        "semantic.timeline.step.1",
        "semantic.metric.value",
        "semantic.process.step.1",
        "semantic.quote.evidence",
        "semantic.table",
    ]

    for theme in ("midnight", "glacier"):
        fixture["theme"] = theme
        exported = Presentation(create_presentation(fixture, tmp_path / f"semantic-blocks-{theme}.pptx"))
        assert len(exported.slides) == 9
        for slide, expected_name in zip(list(exported.slides)[1:], expected_names, strict=True):
            names = {shape.name for shape in slide.shapes}
            assert expected_name in names
            assert all(shape.left >= 0 and shape.top >= 0 for shape in slide.shapes)
            assert all(shape.left + shape.width <= exported.slide_width for shape in slide.shapes)
            assert all(shape.top + shape.height <= exported.slide_height for shape in slide.shapes)
        assert any(shape.has_table for shape in exported.slides[-1].shapes)

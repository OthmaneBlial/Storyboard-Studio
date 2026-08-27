import hashlib
import json
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pydantic import ValidationError

from generate_pptx import create_presentation
from schemas import ChartBlock, LocalAsset, PresentationPayload
from storyboard_studio.assets import AssetValidationError, chart_series, resolve_assets


def digest(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def local_asset(path: Path, *, kind: str, media_type: str, **overrides) -> LocalAsset:
    values = {
        "id": path.stem,
        "kind": kind,
        "path": path.name,
        "sha256": digest(path),
        "media_type": media_type,
        "license": "CC0-1.0",
        "attribution": "Synthetic test fixture",
        "alt_text": "Synthetic image for renderer testing" if kind == "image" else "",
        "source_note": "Synthetic values for renderer testing" if kind == "data" else "",
    }
    values.update(overrides)
    return LocalAsset.model_validate(values)


def test_png_jpeg_svg_csv_and_json_assets_are_local_bounded_and_decodable(tmp_path: Path):
    png = tmp_path / "visual.png"
    jpeg = tmp_path / "photo.jpg"
    svg = tmp_path / "diagram.svg"
    csv_path = tmp_path / "values.csv"
    json_path = tmp_path / "json-values.json"
    Image.new("RGB", (320, 180), "#0a7c86").save(png)
    Image.new("RGB", (320, 180), "#e5b560").save(jpeg)
    svg.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="320" height="180">'
        '<rect width="320" height="180" fill="#123544"/></svg>',
        encoding="utf-8",
    )
    csv_path.write_text("week,value\nOne,2\nTwo,4\n", encoding="utf-8")
    json_path.write_text(json.dumps([{"week": "One", "value": 2}]), encoding="utf-8")
    entries = [
        local_asset(png, kind="image", media_type="image/png"),
        local_asset(jpeg, kind="image", media_type="image/jpeg"),
        local_asset(svg, kind="image", media_type="image/svg+xml"),
        local_asset(csv_path, kind="data", media_type="text/csv"),
        local_asset(json_path, kind="data", media_type="application/json"),
    ]

    resolved = resolve_assets(entries, tmp_path, tmp_path / "cache")

    assert set(resolved) == {"visual", "photo", "diagram", "values", "json-values"}
    assert resolved["diagram"].render_path.suffix == ".png"
    block = ChartBlock(
        chart_type="bar",
        asset_id="values",
        category_field="week",
        value_fields=["value"],
        title="Fixture",
        source_note="Synthetic fixture",
    )
    categories, series = chart_series(resolved["values"], block)
    assert categories == ["One", "Two"]
    assert series == [("value", [2.0, 4.0])]

    presentation = PresentationPayload.model_validate(
        {
            "title": "Local image formats",
            "assets": [entry.model_dump(mode="json") for entry in entries[:3]],
            "slides": [
                {
                    "slide_number": index,
                    "title": f"Local {entry.media_type}",
                    "content": "A checksum-verified local image fixture.",
                    "block": "image",
                    "content_block": {
                        "type": "image",
                        "asset_id": entry.id,
                        "alt_text": entry.alt_text,
                        "fit": "contain",
                    },
                }
                for index, entry in enumerate(entries[:3], start=1)
            ],
        }
    )
    exported = Presentation(
        create_presentation(
            presentation.model_dump(mode="json"),
            tmp_path / "local-image-formats.pptx",
            asset_root=tmp_path,
        )
    )
    assert all(
        any(shape.name.startswith("semantic.image.") for shape in exported.slides[index].shapes)
        for index in (1, 2, 3)
    )


def test_asset_contract_rejects_remote_missing_mismatched_and_active_content(tmp_path: Path):
    with pytest.raises(ValidationError, match="local"):
        LocalAsset.model_validate(
            {
                "id": "remote",
                "kind": "image",
                "path": "https://example.com/image.png",
                "sha256": "0" * 64,
                "media_type": "image/png",
                "license": "CC0",
                "attribution": "Fixture",
                "alt_text": "Remote image",
            }
        )

    with pytest.raises(ValidationError, match="license"):
        LocalAsset.model_validate(
            {
                "id": "unlicensed",
                "kind": "image",
                "path": "image.png",
                "sha256": "0" * 64,
                "media_type": "image/png",
                "license": "",
                "attribution": "Fixture",
                "alt_text": "Synthetic image",
            }
        )

    missing = LocalAsset(
        id="missing",
        kind="data",
        path="missing.csv",
        sha256="0" * 64,
        media_type="text/csv",
        license="CC0",
        attribution="Fixture",
        source_note="Synthetic",
    )
    with pytest.raises(AssetValidationError, match="missing"):
        resolve_assets([missing], tmp_path, tmp_path / "cache")

    csv_path = tmp_path / "values.csv"
    csv_path.write_text("label,value\nA,1\n", encoding="utf-8")
    mismatch = local_asset(csv_path, kind="data", media_type="text/csv").model_copy(
        update={"sha256": "0" * 64}
    )
    with pytest.raises(AssetValidationError, match="checksum mismatch"):
        resolve_assets([mismatch], tmp_path, tmp_path / "cache")

    oversized_csv = tmp_path / "oversized.csv"
    oversized_csv.write_bytes(b"label,value\n" + b"A,1\n" * 70_000)
    oversized = local_asset(oversized_csv, kind="data", media_type="text/csv")
    with pytest.raises(AssetValidationError, match="oversized"):
        resolve_assets([oversized], tmp_path, tmp_path / "cache")

    unreadable_png = tmp_path / "unreadable.png"
    unreadable_png.write_bytes(b"not a png")
    unreadable = local_asset(unreadable_png, kind="image", media_type="image/png")
    with pytest.raises(AssetValidationError, match="unreadable"):
        resolve_assets([unreadable], tmp_path, tmp_path / "cache")

    active_svg = tmp_path / "active.svg"
    active_svg.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg"><image href="https://example.com/a.png"/></svg>',
        encoding="utf-8",
    )
    active = local_asset(active_svg, kind="image", media_type="image/svg+xml")
    with pytest.raises(AssetValidationError, match="non-local href"):
        resolve_assets([active], tmp_path, tmp_path / "cache")


def test_native_visual_fixture_renders_editable_charts_image_alt_text_and_provenance(
    tmp_path: Path,
):
    fixture = Path("assets/demo/native-visuals.json")
    payload = PresentationPayload.model_validate_json(fixture.read_text(encoding="utf-8"))
    output = create_presentation(
        payload.model_dump(mode="json"),
        tmp_path / "native-visuals.pptx",
        asset_root=fixture.parent,
    )
    exported = Presentation(output)

    assert len(exported.slides) == 5
    assert all(any(shape.has_chart for shape in exported.slides[index].shapes) for index in (1, 2, 3))
    picture = next(
        shape for shape in exported.slides[4].shapes if shape.name == "semantic.image.decision-flow"
    )
    assert "local brief moves through diagnosis" in picture._element.nvPicPr.cNvPr.get("descr").lower()
    assert "checksum-verified" in exported.slides[4].notes_slide.notes_text_frame.text

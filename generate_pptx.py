"""A small, dependable PowerPoint renderer for Storyboard Studio.

The renderer has no network dependency. It accepts normalized presentation data
and produces an editable 16:9 PPTX using only native PowerPoint elements.
"""

from __future__ import annotations

import argparse
import json
import tempfile
from collections.abc import Mapping
from importlib.metadata import PackageNotFoundError, version
from pathlib import Path
from typing import Any

from PIL import Image as PILImage
from PIL import ImageOps
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from schemas import ChartBlock, LocalAsset
from storyboard_studio.assets import ResolvedAsset, chart_series, resolve_assets
from storyboard_studio.semantic import normalize_content_block

THEMES: dict[str, dict[str, str]] = {
    "midnight": {
        "name": "Midnight editorial",
        "bg": "101425",
        "surface": "1B2136",
        "surface_alt": "262E49",
        "text": "F7F4EE",
        "muted": "B8C0D6",
        "accent": "E5B560",
        "accent_soft": "745B36",
    },
    "glacier": {
        "name": "Glacier field notes",
        "bg": "F4F8F8",
        "surface": "E4EFF0",
        "surface_alt": "D4E5E7",
        "text": "123544",
        "muted": "55727A",
        "accent": "0A7C86",
        "accent_soft": "B8D9D7",
    },
    "ember": {
        "name": "Ember study",
        "bg": "25120F",
        "surface": "38201A",
        "surface_alt": "4B2A21",
        "text": "FFF5E7",
        "muted": "D6B9A8",
        "accent": "F08A4B",
        "accent_soft": "79432B",
    },
    "forest": {
        "name": "Forest fieldwork",
        "bg": "F1F5EE",
        "surface": "E0EADD",
        "surface_alt": "D0DFC9",
        "text": "1D3829",
        "muted": "577060",
        "accent": "438360",
        "accent_soft": "B9D1BE",
    },
    "royal": {
        "name": "Royal archive",
        "bg": "17151B",
        "surface": "28242D",
        "surface_alt": "37313D",
        "text": "F5ECD8",
        "muted": "C7B99A",
        "accent": "C79C54",
        "accent_soft": "6E5735",
    },
    "sakura": {
        "name": "Sakura notebook",
        "bg": "FCF4F4",
        "surface": "F7E5E7",
        "surface_alt": "F0D5DA",
        "text": "4D2736",
        "muted": "846274",
        "accent": "BE526D",
        "accent_soft": "E5B8C4",
    },
}

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
DISPLAY_FONT = "Aptos Display"
BODY_FONT = "Aptos"

try:
    PACKAGE_VERSION = version("storyboard-studio")
except PackageNotFoundError:  # pragma: no cover - source checkout without metadata
    PACKAGE_VERSION = "0.1.1"


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.replace("#", "").upper())


def _as_text(value: Any, limit: int, fallback: str = "") -> str:
    if not isinstance(value, str):
        return fallback
    value = " ".join(value.split())
    return value[:limit] if value else fallback


def _set_fill(shape: Any, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_text(
    slide: Any,
    text: str,
    x: Any,
    y: Any,
    w: Any,
    h: Any,
    *,
    size: float,
    color: RGBColor,
    font: str = BODY_FONT,
    bold: bool = False,
    align: Any = PP_ALIGN.LEFT,
    valign: Any = MSO_ANCHOR.TOP,
    margins: float = 0.0,
) -> Any:
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margins)
    frame.margin_right = Inches(margins)
    frame.margin_top = Inches(margins)
    frame.margin_bottom = Inches(margins)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _base_slide(prs: Presentation, theme: Mapping[str, str]) -> Any:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(theme["bg"])
    return slide


def _add_footer(slide: Any, data: Mapping[str, Any], theme: Mapping[str, str], page: int) -> None:
    _add_text(
        slide,
        _as_text(data.get("title"), 58),
        Inches(0.55),
        Inches(7.05),
        Inches(6.2),
        Inches(0.22),
        size=9,
        color=_rgb(theme["muted"]),
    )
    _add_text(
        slide,
        str(page).zfill(2),
        Inches(12.1),
        Inches(7.03),
        Inches(0.65),
        Inches(0.22),
        size=9,
        color=_rgb(theme["muted"]),
        align=PP_ALIGN.RIGHT,
    )


def _add_notes(
    slide: Any,
    slide_data: Mapping[str, Any],
    assets: Mapping[str, ResolvedAsset],
) -> None:
    """Write optional author notes and source references to native PPTX notes."""
    notes = _as_text(slide_data.get("speaker_notes"), 1200)
    sources = slide_data.get("sources")
    rows = []
    if isinstance(sources, list):
        for source in sources[:6]:
            if isinstance(source, Mapping):
                label = _as_text(source.get("label"), 100)
                evidence = _as_text(source.get("evidence"), 300)
                owner = _as_text(source.get("owner"), 80)
                if label:
                    rows.append(" — ".join(part for part in (label, evidence, owner) if part))
    if rows:
        source_text = "Sources / evidence (author-supplied; not verified):\n" + "\n".join(rows)
        notes = f"{notes}\n\n{source_text}" if notes else source_text
    block = normalize_content_block(slide_data)
    asset_id = _as_text(block.get("asset_id"), 64)
    if asset_id and asset_id in assets:
        asset = assets[asset_id].entry
        provenance = (
            "Local asset (checksum-verified): "
            f"{asset.id} — {asset.path} — sha256 {asset.sha256} — "
            f"license {asset.license} — attribution {asset.attribution}"
        )
        notes = f"{notes}\n\n{provenance}" if notes else provenance
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _name(shape: Any, value: str) -> Any:
    shape.name = value
    return shape


def _render_standard_block(
    slide: Any,
    block: Mapping[str, Any],
    x: Any,
    width: Any,
    theme: Mapping[str, str],
) -> None:
    points = block.get("points") if isinstance(block.get("points"), list) else []
    row_height = 0.72 if len(points) > 3 else 0.84
    gap = 0.16
    for index, point in enumerate(points[:4]):
        point = point if isinstance(point, Mapping) else {}
        y = Inches(3.05 + index * (row_height + gap))
        row = _name(
            slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, Inches(row_height)),
            f"semantic.standard.point.{index + 1}",
        )
        _set_fill(row, _rgb(theme["surface"] if index % 2 == 0 else theme["surface_alt"]))
        _add_text(
            slide,
            _as_text(point.get("label"), 8, str(index + 1).zfill(2)),
            x + Inches(0.18),
            y + Inches(0.18),
            Inches(0.52),
            Inches(0.26),
            size=11,
            color=_rgb(theme["accent"]),
            bold=True,
        )
        _add_text(
            slide,
            _as_text(point.get("title"), 62, f"Point {index + 1}"),
            x + Inches(0.76),
            y + Inches(0.12),
            Inches(2.25),
            Inches(0.3),
            size=15,
            color=_rgb(theme["text"]),
            bold=True,
        )
        _add_text(
            slide,
            _as_text(point.get("description"), 120),
            x + Inches(3.05),
            y + Inches(0.12),
            width - Inches(3.24),
            Inches(0.4),
            size=13,
            color=_rgb(theme["muted"]),
        )


def _render_comparison_block(
    slide: Any,
    block: Mapping[str, Any],
    x: Any,
    width: Any,
    theme: Mapping[str, str],
) -> None:
    sides = block.get("sides") if isinstance(block.get("sides"), list) else []
    column_width = (width - Inches(0.22)) / 2
    for index, side in enumerate(sides[:2]):
        side = side if isinstance(side, Mapping) else {}
        side_x = x + index * (column_width + Inches(0.22))
        panel = _name(
            slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, side_x, Inches(3.0), column_width, Inches(1.25)),
            f"semantic.comparison.side.{index + 1}",
        )
        _set_fill(panel, _rgb(theme["surface"] if index == 0 else theme["surface_alt"]))
        _add_text(
            slide,
            _as_text(side.get("title"), 70, f"Option {index + 1}"),
            side_x + Inches(0.22),
            Inches(3.18),
            column_width - Inches(0.44),
            Inches(0.32),
            size=17,
            color=_rgb(theme["text"]),
            bold=True,
        )
        _add_text(
            slide,
            _as_text(side.get("summary"), 180),
            side_x + Inches(0.22),
            Inches(3.56),
            column_width - Inches(0.44),
            Inches(0.52),
            size=12,
            color=_rgb(theme["muted"]),
        )
    criteria = block.get("criteria") if isinstance(block.get("criteria"), list) else []
    for index, criterion in enumerate(criteria[:3]):
        criterion = criterion if isinstance(criterion, Mapping) else {}
        y = Inches(4.48 + index * 0.62)
        label = _name(
            slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(1.55), Inches(0.48)),
            f"semantic.comparison.criterion.{index + 1}",
        )
        _set_fill(label, _rgb(theme["accent"] if index == 0 else theme["accent_soft"]))
        _add_text(
            slide,
            _as_text(criterion.get("label"), 60, "Criterion"),
            x + Inches(0.12),
            y + Inches(0.1),
            Inches(1.3),
            Inches(0.25),
            size=10,
            color=_rgb(theme["text"]),
            bold=True,
        )
        remaining = width - Inches(1.75)
        _add_text(
            slide,
            _as_text(criterion.get("left"), 120),
            x + Inches(1.75),
            y + Inches(0.08),
            remaining / 2 - Inches(0.08),
            Inches(0.3),
            size=10,
            color=_rgb(theme["muted"]),
        )
        _add_text(
            slide,
            _as_text(criterion.get("right"), 120),
            x + Inches(1.75) + remaining / 2,
            y + Inches(0.08),
            remaining / 2,
            Inches(0.3),
            size=10,
            color=_rgb(theme["muted"]),
        )


def _render_decision_block(
    slide: Any,
    block: Mapping[str, Any],
    x: Any,
    width: Any,
    theme: Mapping[str, str],
) -> None:
    callout = _name(
        slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x, Inches(3.0), width, Inches(0.8)),
        "semantic.decision.statement",
    )
    _set_fill(callout, _rgb(theme["accent"]))
    _add_text(
        slide,
        _as_text(block.get("decision"), 180, "Decision not supplied"),
        x + Inches(0.24),
        Inches(3.2),
        width - Inches(0.7),
        Inches(0.36),
        size=16,
        color=_rgb(theme["bg"]),
        bold=True,
    )
    options = block.get("options") if isinstance(block.get("options"), list) else []
    option_width = (width - Inches(0.24) * max(0, len(options) - 1)) / max(1, len(options))
    for index, option in enumerate(options[:3]):
        option = option if isinstance(option, Mapping) else {}
        option_x = x + index * (option_width + Inches(0.24))
        panel = _name(
            slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, option_x, Inches(4.08), option_width, Inches(1.12)
            ),
            f"semantic.decision.option.{index + 1}",
        )
        _set_fill(panel, _rgb(theme["surface"] if index % 2 == 0 else theme["surface_alt"]))
        _add_text(
            slide,
            _as_text(option.get("title"), 70, f"Option {index + 1}"),
            option_x + Inches(0.16),
            Inches(4.24),
            option_width - Inches(0.32),
            Inches(0.28),
            size=13,
            color=_rgb(theme["text"]),
            bold=True,
        )
        _add_text(
            slide,
            _as_text(option.get("description"), 220),
            option_x + Inches(0.16),
            Inches(4.58),
            option_width - Inches(0.32),
            Inches(0.42),
            size=9,
            color=_rgb(theme["muted"]),
        )
    rationale = _name(
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(5.48), width, Inches(0.76)),
        "semantic.decision.rationale",
    )
    _set_fill(rationale, _rgb(theme["surface_alt"]))
    rationale_text = _as_text(block.get("rationale"), 220)
    owner = _as_text(block.get("owner"), 80)
    _add_text(
        slide,
        f"RATIONALE  {rationale_text}" + (f"  ·  OWNER  {owner}" if owner else ""),
        x + Inches(0.2),
        Inches(5.67),
        width - Inches(0.4),
        Inches(0.32),
        size=11,
        color=_rgb(theme["text"]),
        bold=True,
    )


def _render_timeline_block(
    slide: Any,
    block: Mapping[str, Any],
    x: Any,
    width: Any,
    theme: Mapping[str, str],
) -> None:
    steps = block.get("steps") if isinstance(block.get("steps"), list) else []
    count = max(1, len(steps))
    lane_width = width / count
    line = _name(
        slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x + Inches(0.35), Inches(3.52), width - Inches(0.7), Inches(0.08)
        ),
        "semantic.timeline.connector",
    )
    _set_fill(line, _rgb(theme["accent_soft"]))
    for index, step in enumerate(steps[:4]):
        step = step if isinstance(step, Mapping) else {}
        lane_x = x + index * lane_width
        dot = _name(
            slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                lane_x + lane_width / 2 - Inches(0.2),
                Inches(3.34),
                Inches(0.4),
                Inches(0.4),
            ),
            f"semantic.timeline.step.{index + 1}",
        )
        _set_fill(dot, _rgb(theme["accent"]))
        _add_text(
            slide,
            _as_text(step.get("label"), 24, str(index + 1)),
            lane_x + Inches(0.08),
            Inches(3.0),
            lane_width - Inches(0.16),
            Inches(0.28),
            size=11,
            color=_rgb(theme["accent"]),
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        _add_text(
            slide,
            _as_text(step.get("title"), 80),
            lane_x + Inches(0.08),
            Inches(4.02),
            lane_width - Inches(0.16),
            Inches(0.78),
            size=14,
            color=_rgb(theme["text"]),
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        _add_text(
            slide,
            _as_text(step.get("owner"), 80),
            lane_x + Inches(0.08),
            Inches(4.98),
            lane_width - Inches(0.16),
            Inches(0.36),
            size=10,
            color=_rgb(theme["muted"]),
            align=PP_ALIGN.CENTER,
        )


def _render_metric_block(
    slide: Any,
    block: Mapping[str, Any],
    x: Any,
    width: Any,
    theme: Mapping[str, str],
) -> None:
    metric = _name(
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(3.0), width * 0.4, Inches(2.75)),
        "semantic.metric.value",
    )
    _set_fill(metric, _rgb(theme["accent"]))
    _add_text(
        slide,
        _as_text(block.get("value"), 24, "—"),
        x + Inches(0.2),
        Inches(3.38),
        width * 0.4 - Inches(0.4),
        Inches(1.02),
        size=46,
        color=_rgb(theme["bg"]),
        font=DISPLAY_FONT,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        _as_text(block.get("label"), 80),
        x + Inches(0.2),
        Inches(4.58),
        width * 0.4 - Inches(0.4),
        Inches(0.5),
        size=15,
        color=_rgb(theme["bg"]),
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    detail_x = x + width * 0.44
    _add_text(
        slide,
        _as_text(block.get("context"), 220),
        detail_x,
        Inches(3.25),
        width * 0.56,
        Inches(1.35),
        size=22,
        color=_rgb(theme["text"]),
        font=DISPLAY_FONT,
        bold=True,
    )
    _add_text(
        slide,
        "SOURCE  " + _as_text(block.get("source"), 120, "Not supplied"),
        detail_x,
        Inches(5.05),
        width * 0.56,
        Inches(0.42),
        size=11,
        color=_rgb(theme["muted"]),
    )


def _render_process_block(
    slide: Any,
    block: Mapping[str, Any],
    x: Any,
    width: Any,
    theme: Mapping[str, str],
) -> None:
    steps = block.get("steps") if isinstance(block.get("steps"), list) else []
    step_width = width / max(1, len(steps))
    for index, step in enumerate(steps[:5]):
        step = step if isinstance(step, Mapping) else {}
        step_x = x + index * step_width
        chevron = _name(
            slide.shapes.add_shape(
                MSO_SHAPE.CHEVRON,
                step_x,
                Inches(3.05 + (index % 2) * 0.28),
                step_width + Inches(0.12),
                Inches(1.0),
            ),
            f"semantic.process.step.{index + 1}",
        )
        _set_fill(chevron, _rgb(theme["accent"] if index == 0 else theme["surface_alt"]))
        _add_text(
            slide,
            _as_text(step.get("title"), 70, f"Step {index + 1}"),
            step_x + Inches(0.12),
            Inches(3.32 + (index % 2) * 0.28),
            step_width - Inches(0.1),
            Inches(0.32),
            size=12,
            color=_rgb(theme["bg"] if index == 0 else theme["text"]),
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        _add_text(
            slide,
            _as_text(step.get("description"), 140),
            step_x + Inches(0.12),
            Inches(4.55),
            step_width - Inches(0.18),
            Inches(1.1),
            size=11,
            color=_rgb(theme["muted"]),
            align=PP_ALIGN.CENTER,
        )


def _render_quote_block(
    slide: Any,
    block: Mapping[str, Any],
    x: Any,
    width: Any,
    theme: Mapping[str, str],
) -> None:
    quote_panel = _name(
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(3.0), width, Inches(2.8)),
        "semantic.quote.evidence",
    )
    _set_fill(quote_panel, _rgb(theme["surface"]))
    _add_text(
        slide,
        "“",
        x + Inches(0.28),
        Inches(3.08),
        Inches(0.7),
        Inches(0.7),
        size=50,
        color=_rgb(theme["accent"]),
        font=DISPLAY_FONT,
        bold=True,
    )
    _add_text(
        slide,
        _as_text(block.get("quote"), 280),
        x + Inches(1.0),
        Inches(3.35),
        width - Inches(1.35),
        Inches(1.18),
        size=23,
        color=_rgb(theme["text"]),
        font=DISPLAY_FONT,
        bold=True,
    )
    _add_text(
        slide,
        "— " + _as_text(block.get("attribution"), 100, "Attribution not supplied"),
        x + Inches(1.0),
        Inches(4.8),
        width - Inches(1.35),
        Inches(0.34),
        size=12,
        color=_rgb(theme["accent"]),
        bold=True,
    )
    _add_text(
        slide,
        "EVIDENCE  " + _as_text(block.get("evidence"), 180, "Not supplied"),
        x + Inches(1.0),
        Inches(5.25),
        width - Inches(1.35),
        Inches(0.34),
        size=10,
        color=_rgb(theme["muted"]),
    )


def _render_table_block(
    slide: Any,
    block: Mapping[str, Any],
    x: Any,
    width: Any,
    theme: Mapping[str, str],
) -> None:
    columns = block.get("columns") if isinstance(block.get("columns"), list) else []
    rows = block.get("rows") if isinstance(block.get("rows"), list) else []
    table_shape = slide.shapes.add_table(
        len(rows) + 1,
        len(columns),
        x,
        Inches(3.0),
        width,
        Inches(min(3.0, 0.56 * (len(rows) + 1))),
    )
    _name(table_shape, "semantic.table")
    table = table_shape.table
    for column_index, heading in enumerate(columns):
        table.columns[column_index].width = int(width / len(columns))
        cell = table.cell(0, column_index)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(theme["accent"])
        cell.text = _as_text(heading, 60)
    for row_index, row in enumerate(rows, start=1):
        cells = row.get("cells") if isinstance(row, Mapping) else []
        for column_index, value in enumerate(cells):
            cell = table.cell(row_index, column_index)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(theme["surface"] if row_index % 2 else theme["surface_alt"])
            cell.text = _as_text(value, 100)
    for row_index in range(len(rows) + 1):
        for column_index in range(len(columns)):
            cell = table.cell(row_index, column_index)
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.07)
            cell.margin_bottom = Inches(0.05)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = BODY_FONT
                    run.font.size = Pt(11 if row_index else 12)
                    run.font.bold = row_index == 0
                    run.font.color.rgb = _rgb(theme["bg"] if row_index == 0 else theme["text"])


def _asset_for(block: Mapping[str, Any], assets: Mapping[str, ResolvedAsset]) -> ResolvedAsset:
    asset_id = _as_text(block.get("asset_id"), 64)
    if not asset_id or asset_id not in assets:
        raise ValueError(f"Semantic block references unavailable local asset {asset_id!r}.")
    return assets[asset_id]


def _render_chart_block(
    slide: Any,
    block: Mapping[str, Any],
    x: Any,
    width: Any,
    theme: Mapping[str, str],
    assets: Mapping[str, ResolvedAsset],
) -> None:
    asset = _asset_for(block, assets)
    contract = ChartBlock.model_validate(block)
    categories, series = chart_series(asset, contract)
    data = ChartData()
    data.categories = categories
    for name, values in series:
        data.add_series(name, values)
    chart_types = {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "donut": XL_CHART_TYPE.DOUGHNUT,
    }
    chart_shape = slide.shapes.add_chart(
        chart_types[contract.chart_type],
        x,
        Inches(3.0),
        width,
        Inches(2.72),
        data,
    )
    _name(chart_shape, f"semantic.chart.{contract.chart_type}")
    chart = chart_shape.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = contract.title
    chart.has_legend = len(series) > 1 or contract.chart_type == "donut"
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    for index, plot in enumerate(chart.plots):
        plot.has_data_labels = True
        plot.data_labels.position = (
            XL_LABEL_POSITION.BEST_FIT if contract.chart_type == "donut" else XL_LABEL_POSITION.OUTSIDE_END
        )
        for chart_series_item in plot.series:
            chart_series_item.format.fill.solid()
            chart_series_item.format.fill.fore_color.rgb = _rgb(
                theme["accent"] if index == 0 else theme["accent_soft"]
            )
    if contract.chart_type != "donut":
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.tick_labels.font.size = Pt(9)
        chart.category_axis.tick_labels.font.size = Pt(9)
    _add_text(
        slide,
        f"SOURCE  {contract.source_note}  ·  {asset.entry.attribution}",
        x,
        Inches(5.88),
        width,
        Inches(0.38),
        size=9,
        color=_rgb(theme["muted"]),
    )


def _cover_image(asset: ResolvedAsset, width: Any, height: Any, cache_dir: Path) -> Path:
    target_width = 1600
    target_height = max(1, round(target_width * int(height) / int(width)))
    destination = cache_dir / f"{asset.entry.id}-cover.png"
    with PILImage.open(asset.render_path) as image:
        fitted = ImageOps.fit(image.convert("RGBA"), (target_width, target_height))
        fitted.save(destination, format="PNG")
    return destination


def _render_image_block(
    slide: Any,
    block: Mapping[str, Any],
    x: Any,
    width: Any,
    theme: Mapping[str, str],
    assets: Mapping[str, ResolvedAsset],
    cache_dir: Path,
) -> None:
    asset = _asset_for(block, assets)
    image_height = Inches(2.72)
    image_path = (
        _cover_image(asset, width, image_height, cache_dir)
        if block.get("fit") == "cover"
        else asset.render_path
    )
    if block.get("fit") == "cover":
        picture = slide.shapes.add_picture(str(image_path), x, Inches(3.0), width=width, height=image_height)
    else:
        source_width = max(1, asset.width or 1)
        source_height = max(1, asset.height or 1)
        scale = min(int(width) / source_width, int(image_height) / source_height)
        rendered_width = int(source_width * scale)
        rendered_height = int(source_height * scale)
        picture = slide.shapes.add_picture(
            str(image_path),
            x + (int(width) - rendered_width) // 2,
            Inches(3.0) + (int(image_height) - rendered_height) // 2,
            width=rendered_width,
            height=rendered_height,
        )
    _name(picture, f"semantic.image.{asset.entry.id}")
    picture._element.nvPicPr.cNvPr.set("descr", _as_text(block.get("alt_text"), 240))
    caption = _as_text(block.get("caption"), 160)
    provenance = f"{asset.entry.attribution}  ·  {asset.entry.license}"
    _add_text(
        slide,
        "  ·  ".join(part for part in (caption, provenance) if part),
        x,
        Inches(5.88),
        width,
        Inches(0.38),
        size=9,
        color=_rgb(theme["muted"]),
    )


def _render_content_block(
    slide: Any,
    block: Mapping[str, Any],
    x: Any,
    width: Any,
    theme: Mapping[str, str],
    assets: Mapping[str, ResolvedAsset],
    cache_dir: Path,
) -> None:
    renderers = {
        "standard": _render_standard_block,
        "comparison": _render_comparison_block,
        "decision": _render_decision_block,
        "timeline": _render_timeline_block,
        "metric": _render_metric_block,
        "process": _render_process_block,
        "quote": _render_quote_block,
        "table": _render_table_block,
    }
    if block.get("type") == "chart":
        _render_chart_block(slide, block, x, width, theme, assets)
        return
    if block.get("type") == "image":
        _render_image_block(slide, block, x, width, theme, assets, cache_dir)
        return
    renderers.get(str(block.get("type")), _render_standard_block)(slide, block, x, width, theme)


def _add_title_slide(prs: Presentation, data: Mapping[str, Any], theme: Mapping[str, str]) -> None:
    slide = _base_slide(prs, theme)
    bg = _rgb(theme["bg"])
    surface = _rgb(theme["surface"])
    accent = _rgb(theme["accent"])
    text = _rgb(theme["text"])
    muted = _rgb(theme["muted"])

    left_rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.65), Inches(0.1), Inches(5.9)
    )
    _set_fill(left_rule, accent)
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.75), Inches(0), Inches(3.583), Inches(7.5))
    _set_fill(panel, surface)
    lower_panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(9.75), Inches(4.7), Inches(3.583), Inches(2.8)
    )
    _set_fill(lower_panel, accent)
    small_block = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(10.35), Inches(1.0), Inches(2.4), Inches(0.16)
    )
    _set_fill(small_block, accent)

    _add_text(
        slide,
        "STORYBOARD STUDIO",
        Inches(1.0),
        Inches(0.84),
        Inches(5.7),
        Inches(0.32),
        size=12,
        color=accent,
        bold=True,
    )
    _add_text(
        slide,
        _as_text(data.get("title"), 90, "Untitled presentation"),
        Inches(0.95),
        Inches(1.48),
        Inches(8.1),
        Inches(2.45),
        size=48,
        color=text,
        font=DISPLAY_FONT,
        bold=True,
    )
    _add_text(
        slide,
        _as_text(data.get("subtitle"), 110),
        Inches(1.0),
        Inches(4.28),
        Inches(7.55),
        Inches(0.85),
        size=21,
        color=muted,
    )
    _add_text(
        slide,
        "A concise, editable deck",
        Inches(10.35),
        Inches(1.42),
        Inches(2.25),
        Inches(0.6),
        size=17,
        color=text,
        bold=True,
    )
    _add_text(
        slide,
        "Prepared for a clear conversation, not a wall of slides.",
        Inches(10.35),
        Inches(2.25),
        Inches(2.25),
        Inches(1.05),
        size=16,
        color=muted,
    )
    _add_text(
        slide,
        "EDITABLE\nPOWERPOINT",
        Inches(10.35),
        Inches(5.25),
        Inches(2.2),
        Inches(1.2),
        size=20,
        color=bg,
        font=DISPLAY_FONT,
        bold=True,
    )


def _add_content_slide(
    prs: Presentation,
    data: Mapping[str, Any],
    slide_data: Mapping[str, Any],
    theme: Mapping[str, str],
    page: int,
    assets: Mapping[str, ResolvedAsset],
    cache_dir: Path,
) -> None:
    slide = _base_slide(prs, theme)
    text = _rgb(theme["text"])
    muted = _rgb(theme["muted"])
    accent = _rgb(theme["accent"])
    surface = _rgb(theme["surface"])
    surface_alt = _rgb(theme["surface_alt"])
    layout = slide_data.get("layout") if slide_data.get("layout") in {"left", "right", "focus"} else "right"
    content_block = normalize_content_block(slide_data)
    block = str(content_block.get("type", "standard"))

    if layout == "focus":
        visual_x, visual_w = Inches(9.45), Inches(3.15)
        content_x, content_w = Inches(0.6), Inches(8.25)
    elif layout == "left":
        visual_x, visual_w = Inches(0.6), Inches(3.5)
        content_x, content_w = Inches(4.55), Inches(7.8)
    else:
        visual_x, visual_w = Inches(9.15), Inches(3.65)
        content_x, content_w = Inches(0.6), Inches(7.8)

    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.13), SLIDE_HEIGHT)
    _set_fill(accent_bar, accent)
    block_labels = {
        "standard": "KEY FRAME",
        "comparison": "COMPARE",
        "decision": "DECISION",
        "timeline": "SEQUENCE",
        "metric": "SIGNAL",
        "process": "PROCESS",
        "quote": "EVIDENCE",
        "table": "TABLE",
        "chart": "DATA",
        "image": "VISUAL",
    }
    _add_text(
        slide,
        "STORYBOARD / " + str(page).zfill(2),
        Inches(0.58),
        Inches(0.42),
        Inches(2.2),
        Inches(0.28),
        size=10,
        color=accent,
        bold=True,
    )
    heading_x = content_x if layout == "left" else Inches(0.55)
    heading_width = content_w if layout == "left" else Inches(7.7 if layout != "focus" else 11.6)
    intro_x = content_x if layout == "left" else Inches(0.6)
    intro_width = content_w if layout == "left" else Inches(7.0 if layout != "focus" else 8.35)

    _add_text(
        slide,
        _as_text(slide_data.get("title"), 68, f"Slide {page}"),
        heading_x,
        Inches(0.88),
        heading_width,
        Inches(0.75),
        size=35,
        color=text,
        font=DISPLAY_FONT,
        bold=True,
    )
    _add_text(
        slide,
        _as_text(slide_data.get("content"), 220),
        intro_x,
        Inches(1.82),
        intro_width,
        Inches(1.12),
        size=18,
        color=muted,
    )

    visual = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, visual_x, Inches(1.78), visual_w, Inches(4.58)
    )
    _set_fill(visual, surface)
    visual_inner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, visual_x + Inches(0.32), Inches(2.12), visual_w - Inches(0.64), Inches(0.12)
    )
    _set_fill(visual_inner, accent)
    _add_text(
        slide,
        str(page).zfill(2),
        visual_x + Inches(0.32),
        Inches(2.5),
        visual_w - Inches(0.64),
        Inches(1.25),
        size=58,
        color=accent,
        font=DISPLAY_FONT,
        bold=True,
    )
    _add_text(
        slide,
        block_labels.get(block, "KEY FRAME"),
        visual_x + Inches(0.34),
        Inches(3.87),
        visual_w - Inches(0.68),
        Inches(0.25),
        size=11,
        color=muted,
        bold=True,
    )
    if block == "comparison":
        split_width = (visual_w - Inches(0.78)) / 2
        before = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, visual_x + Inches(0.34), Inches(5.15), split_width, Inches(0.55)
        )
        _set_fill(before, accent)
        _add_text(
            slide,
            "BEFORE",
            before.left,
            before.top + Inches(0.12),
            before.width,
            Inches(0.25),
            size=9,
            color=text,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        after = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            visual_x + Inches(0.44) + split_width,
            Inches(5.15),
            split_width,
            Inches(0.55),
        )
        _set_fill(after, surface_alt)
        _add_text(
            slide,
            "AFTER",
            after.left,
            after.top + Inches(0.12),
            after.width,
            Inches(0.25),
            size=9,
            color=text,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    elif block == "timeline":
        for marker in range(3):
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                visual_x + Inches(0.5) + Inches(marker * 0.82),
                Inches(5.15),
                Inches(0.28),
                Inches(0.28),
            )
            _set_fill(dot, accent if marker == 0 else surface_alt)
    elif block == "metric":
        _add_text(
            slide,
            _as_text(content_block.get("value"), 24, "—"),
            visual_x + Inches(0.34),
            Inches(5.0),
            visual_w - Inches(0.68),
            Inches(0.72),
            size=32,
            color=accent,
            font=DISPLAY_FONT,
            bold=True,
        )
    elif block == "decision":
        decision = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON, visual_x + Inches(0.36), Inches(5.14), visual_w - Inches(0.72), Inches(0.52)
        )
        _set_fill(decision, accent)
    elif block == "process":
        for marker in range(3):
            chevron = slide.shapes.add_shape(
                MSO_SHAPE.CHEVRON,
                visual_x + Inches(0.36 + marker * 0.75),
                Inches(5.14),
                Inches(0.9),
                Inches(0.5),
            )
            _set_fill(chevron, accent if marker == 0 else surface_alt)
    elif block == "quote":
        _add_text(
            slide,
            "“ ”",
            visual_x + Inches(0.34),
            Inches(4.98),
            visual_w - Inches(0.68),
            Inches(0.7),
            size=36,
            color=accent,
            font=DISPLAY_FONT,
            bold=True,
        )
    elif block == "table":
        for row_index in range(3):
            row = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                visual_x + Inches(0.36),
                Inches(5.02 + row_index * 0.2),
                visual_w - Inches(0.72),
                Inches(0.12),
            )
            _set_fill(row, accent if row_index == 0 else surface_alt)
    elif block == "chart":
        for marker, height in enumerate((0.28, 0.52, 0.78)):
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                visual_x + Inches(0.42 + marker * 0.62),
                Inches(5.74 - height),
                Inches(0.38),
                Inches(height),
            )
            _set_fill(bar, accent if marker == 2 else surface_alt)
    elif block == "image":
        frame = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            visual_x + Inches(0.38),
            Inches(4.98),
            visual_w - Inches(0.76),
            Inches(0.68),
        )
        frame.fill.background()
        frame.line.color.rgb = accent
    _add_text(
        slide,
        _as_text(slide_data.get("title"), 48),
        visual_x + Inches(0.34),
        Inches(4.28),
        visual_w - Inches(0.68),
        Inches(1.12),
        size=19,
        color=text,
        font=DISPLAY_FONT,
        bold=True,
    )

    _render_content_block(slide, content_block, content_x, content_w, theme, assets, cache_dir)
    _add_footer(slide, data, theme, page)
    _add_notes(slide, slide_data, assets)


def create_presentation(
    data: Mapping[str, Any],
    output_path: str | Path = "output/storyboard-presentation.pptx",
    *,
    provenance: str = "",
    asset_root: str | Path | None = None,
) -> Path:
    """Render a validated payload to a new file and return its absolute path."""
    temporary = tempfile.TemporaryDirectory(prefix="storyboard-assets-")
    try:
        cache_dir = Path(temporary.name)
        raw_assets = data.get("assets") if isinstance(data.get("assets"), list) else []
        asset_models = [LocalAsset.model_validate(item) for item in raw_assets]
        assets = resolve_assets(
            asset_models,
            Path(asset_root).expanduser().resolve() if asset_root else Path.cwd(),
            cache_dir,
        )
        theme_id = data.get("theme", "midnight") if isinstance(data, Mapping) else "midnight"
        theme = THEMES.get(str(theme_id), THEMES["midnight"])
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        prs.core_properties.title = _as_text(data.get("title"), 90, "Storyboard Studio presentation")
        prs.core_properties.subject = _as_text(data.get("subtitle"), 110)
        prs.core_properties.author = "Storyboard Studio"
        prs.core_properties.comments = _as_text(provenance, 500)

        _add_title_slide(prs, data, theme)
        slides = data.get("slides") if isinstance(data.get("slides"), list) else []
        for page, slide_data in enumerate(slides, start=1):
            if isinstance(slide_data, Mapping):
                _add_content_slide(prs, data, slide_data, theme, page, assets, cache_dir)

        destination = Path(output_path).expanduser().resolve()
        destination.parent.mkdir(parents=True, exist_ok=True)
        prs.save(destination)
        return destination
    finally:
        temporary.cleanup()


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Render Storyboard Studio JSON into an editable PowerPoint deck."
    )
    parser.add_argument("--version", action="version", version=f"%(prog)s {PACKAGE_VERSION}")
    parser.add_argument("--input", required=True, type=Path, help="Validated presentation JSON file")
    parser.add_argument(
        "--output", type=Path, default=Path("output/storyboard.pptx"), help="Destination PPTX path"
    )
    args = parser.parse_args()
    try:
        with args.input.open(encoding="utf-8") as file:
            from schemas import PresentationPayload

            data = PresentationPayload.model_validate(json.load(file)).model_dump()
        output = create_presentation(data, args.output, asset_root=args.input.parent)
    except (OSError, ValueError, json.JSONDecodeError) as exc:
        parser.error(f"Could not create the presentation: {exc}")
    print(f"Created {output}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

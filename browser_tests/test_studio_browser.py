from __future__ import annotations

import json
import os
import socket
import subprocess
import sys
import time
import zipfile
from pathlib import Path
from urllib.error import URLError
from urllib.request import urlopen

import pytest
from playwright.sync_api import Page, expect, sync_playwright
from pptx import Presentation


def _free_port() -> int:
    with socket.socket() as sock:
        sock.bind(("127.0.0.1", 0))
        return int(sock.getsockname()[1])


@pytest.fixture(scope="module")
def studio_url(tmp_path_factory: pytest.TempPathFactory):
    work = tmp_path_factory.mktemp("browser-studio")
    port = _free_port()
    env = {
        **os.environ,
        "GEMINI_API_KEY": "",
        "STORYBOARD_OUTPUT_DIR": str(work / "output"),
    }
    process = subprocess.Popen(
        [
            sys.executable,
            "-m",
            "storyboard_studio.cli",
            "serve",
            "--host",
            "127.0.0.1",
            "--port",
            str(port),
        ],
        cwd=work,
        env=env,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )
    url = f"http://127.0.0.1:{port}"
    deadline = time.monotonic() + 15
    while time.monotonic() < deadline:
        try:
            with urlopen(f"{url}/api/health", timeout=2) as response:
                if response.status == 200:
                    break
        except (URLError, OSError):
            time.sleep(0.1)
    else:
        process.terminate()
        raise RuntimeError("Packaged Storyboard Studio did not become ready")
    yield url
    process.terminate()
    try:
        process.wait(timeout=5)
    except subprocess.TimeoutExpired:
        process.kill()


@pytest.fixture(scope="module")
def asset_studio_url(tmp_path_factory: pytest.TempPathFactory):
    work = tmp_path_factory.mktemp("browser-assets")
    port = _free_port()
    env = {
        **os.environ,
        "GEMINI_API_KEY": "",
        "STORYBOARD_OUTPUT_DIR": str(work / "output"),
    }
    process = subprocess.Popen(
        [
            sys.executable,
            "-m",
            "storyboard_studio.cli",
            "serve",
            "--host",
            "127.0.0.1",
            "--port",
            str(port),
        ],
        cwd=Path("assets/demo").resolve(),
        env=env,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )
    url = f"http://127.0.0.1:{port}"
    deadline = time.monotonic() + 15
    while time.monotonic() < deadline:
        try:
            with urlopen(f"{url}/api/health", timeout=2) as response:
                if response.status == 200:
                    break
        except (URLError, OSError):
            time.sleep(0.1)
    else:
        process.terminate()
        raise RuntimeError("Asset-aware Storyboard Studio did not become ready")
    yield url
    process.terminate()
    try:
        process.wait(timeout=5)
    except subprocess.TimeoutExpired:
        process.kill()


def _assert_no_horizontal_overflow(page: Page) -> None:
    dimensions = page.evaluate(
        """() => ({
          viewport: window.innerWidth,
          root: document.documentElement.scrollWidth,
          body: document.body.scrollWidth,
        })"""
    )
    assert dimensions["root"] <= dimensions["viewport"]
    assert dimensions["body"] <= dimensions["viewport"]


def test_keyboard_authoring_export_and_responsive_contract(studio_url: str, tmp_path: Path):
    with sync_playwright() as playwright:
        browser = playwright.chromium.launch()
        page = browser.new_page(viewport={"width": 1440, "height": 1000})
        page.goto(studio_url, wait_until="domcontentloaded")

        page.keyboard.press("Tab")
        assert page.locator(":focus").inner_text() == "Skip to the presentation brief"

        sample = page.get_by_role("button", name="Try a sample brief")
        sample.focus()
        page.keyboard.press("Enter")
        assert "onboarding pilot" in page.get_by_label("Decision to make").input_value().lower()

        build = page.get_by_role("button", name="Build decision story")
        build.focus()
        page.keyboard.press("Enter")
        page.locator("#previewSection").wait_for(state="visible")
        assert page.locator("#previewSource").inner_text() == "LOCAL DECISION STORY"
        assert page.locator("#storyMap li").count() == 5

        page.get_by_role("button", name="Run Narrative Doctor").click()
        expect(page.locator("#doctorSummary")).not_to_contain_text("No diagnosis yet")
        first_findings = page.locator(".doctor-finding").count()
        assert first_findings >= 2
        page.get_by_role("button", name="Accept action").first.click()
        assert page.locator(":focus").get_attribute("aria-label") == "Evidence owner for slide 1"
        page.get_by_label("Evidence owner for slide 1").fill("Customer success lead")
        page.get_by_role("button", name="Run Narrative Doctor").click()
        expect(page.locator("#doctorSummary")).to_contain_text("1 notes")
        assert page.locator(".doctor-finding").count() < first_findings

        title = page.get_by_label("Presentation title")
        edited_title = "A practical plan to make remote onboarding feel human and measurable"
        title.fill(edited_title)
        title.press("Tab")
        assert "Unsaved storyboard edits" in page.locator("#saveStatus").inner_text()

        first_slide_title = page.get_by_label("Slide 1 title").input_value()
        page.get_by_role("button", name="Move slide down").first.click()
        assert page.get_by_label("Slide 2 title").input_value() == first_slide_title
        page.get_by_role("button", name="Undo").click()
        assert page.get_by_label("Slide 1 title").input_value() == first_slide_title
        page.get_by_role("button", name="Redo").click()
        assert page.get_by_label("Slide 2 title").input_value() == first_slide_title

        slide_total = page.locator("#storyMap li").count()
        preview_copy = [
            {
                "title": page.get_by_label(f"Slide {index} title").input_value(),
                "summary": page.get_by_label(f"Slide {index} summary").input_value(),
                "semantic_parts": page.locator(".semantic-fallback p")
                .nth(index - 1)
                .inner_text()
                .split(" | "),
            }
            for index in range(1, slide_total + 1)
        ]
        role_labels = {
            "standard": "KEY FRAME",
            "comparison": "COMPARE",
            "decision": "DECISION",
            "timeline": "SEQUENCE",
            "metric": "SIGNAL",
            "process": "PROCESS",
            "quote": "EVIDENCE",
            "table": "TABLE",
        }
        preview_roles = [
            role_labels[
                page.locator("#deckPreview").get_by_label(f"Editorial block for slide {index}").input_value()
            ]
            for index in range(1, slide_total + 1)
        ]

        with page.expect_download() as download_event:
            page.get_by_role("button", name="Export PowerPoint").click()
        downloaded = tmp_path / "storyboard-browser-contract.pptx"
        download_event.value.save_as(downloaded)
        assert downloaded.read_bytes()[:2] == b"PK"
        exported = Presentation(downloaded)
        exported_text = [
            "\n".join(shape.text for shape in slide.shapes if shape.has_text_frame)
            for slide in exported.slides
        ]
        assert edited_title in exported_text[0]
        assert first_slide_title in exported_text[2]
        for slide_number, copy in enumerate(preview_copy, start=1):
            assert copy["title"] in exported_text[slide_number]
            assert copy["summary"] in exported_text[slide_number]
            assert all(part in exported_text[slide_number] for part in copy["semantic_parts"])
            assert preview_roles[slide_number - 1] in exported_text[slide_number]

        with page.expect_download() as bundle_event:
            page.get_by_role("button", name="Export review bundle").click()
        bundle = tmp_path / "storyboard-review-bundle.zip"
        bundle_event.value.save_as(bundle)
        with zipfile.ZipFile(bundle) as archive:
            assert set(archive.namelist()) == {
                "deck.pptx",
                "deck.receipt.json",
                "deck.story.json",
            }

        for width in (320, 375, 1440):
            page.set_viewport_size({"width": width, "height": 900})
            _assert_no_horizontal_overflow(page)
            assert page.get_by_label("Presentation title").is_visible()
            assert page.get_by_role("button", name="Export PowerPoint").is_visible()
            more_actions = page.locator(".preview-more-actions summary")
            if width < 560:
                assert more_actions.is_visible()
                assert not page.get_by_role("button", name="Undo").is_visible()
                more_actions.click()
                assert page.get_by_role("button", name="Undo").is_visible()
                more_actions.click()
            else:
                assert not more_actions.is_visible()
                assert page.get_by_role("button", name="Undo").is_visible()

        title_metrics = title.evaluate(
            "element => ({scrollHeight: element.scrollHeight, clientHeight: element.clientHeight})"
        )
        assert title_metrics["scrollHeight"] <= title_metrics["clientHeight"] + 1
        assert not page.locator("#previewSection [aria-label]:not([type=hidden])").evaluate_all(
            "elements => elements.some(element => element.getClientRects().length === 0)"
        )
        browser.close()


def test_all_semantic_blocks_have_specific_controls_and_plain_text_fallback(studio_url: str, tmp_path: Path):
    fixture = Path("examples/fixtures/semantic-blocks.json").resolve()
    with sync_playwright() as playwright:
        browser = playwright.chromium.launch()
        page = browser.new_page(viewport={"width": 1440, "height": 1000})
        page.goto(studio_url, wait_until="domcontentloaded")
        page.locator("#importOutlineInput").set_input_files(fixture)
        page.locator("#previewSection").wait_for(state="visible")

        assert page.locator(".semantic-editor").count() == 8
        assert page.locator(".semantic-fallback p").count() == 8
        for label in (
            "Slide 1 point 1 title",
            "Slide 2 side 1 title",
            "Slide 3 decision statement",
            "Slide 4 step 1 owner",
            "Slide 5 metric value",
            "Slide 6 process step 1 title",
            "Slide 7 quote attribution",
            "Slide 8 row 1 cell 1",
            "Slide 8 table summary",
        ):
            assert page.get_by_label(label).is_visible()
        assert all(
            (page.locator(".semantic-fallback p").nth(index).text_content() or "").strip()
            for index in range(8)
        )

        metric = page.get_by_label("Slide 5 metric value")
        metric.fill("8/8")
        metric.press("Tab")
        expect(page.locator(".semantic-fallback p").nth(4)).to_contain_text("8/8")

        with page.expect_download() as download_event:
            page.get_by_role("button", name="Export PowerPoint").click()
        downloaded = tmp_path / "semantic-blocks.pptx"
        download_event.value.save_as(downloaded)
        exported = Presentation(downloaded)
        exported_text = "\n".join(shape.text for shape in exported.slides[5].shapes if shape.has_text_frame)
        assert "8/8" in exported_text

        page.locator("#deckPreview").get_by_label("Editorial block for slide 1").select_option("quote")
        assert page.get_by_label("Slide 1 quote", exact=True).is_visible()
        assert page.get_by_label("Slide 1 quote attribution").is_visible()
        assert (
            page.locator(".semantic-editor").first.get_attribute("aria-label")
            == "quote block fields for slide 1"
        )
        browser.close()


def test_local_chart_and_image_assets_show_provenance_and_export(asset_studio_url: str, tmp_path: Path):
    fixture = Path("assets/demo/native-visuals.json").resolve()
    with sync_playwright() as playwright:
        browser = playwright.chromium.launch()
        page = browser.new_page(viewport={"width": 1440, "height": 1000})
        page.goto(asset_studio_url, wait_until="domcontentloaded")
        page.locator("#importOutlineInput").set_input_files(fixture)
        page.locator("#previewSection").wait_for(state="visible")

        assert page.locator(".semantic-chart").count() == 3
        assert page.locator(".semantic-image").count() == 1
        assert page.get_by_label("Slide 1 chart type").is_visible()
        assert page.get_by_label("Slide 1 chart asset").input_value() == "pilot-results"
        assert page.get_by_label("Slide 4 image alt text").is_visible()
        assert page.locator(".asset-provenance").count() == 4
        assert all(
            "SHA-256" in (page.locator(".asset-provenance").nth(index).text_content() or "")
            for index in range(4)
        )

        with page.expect_download() as download_event:
            page.get_by_role("button", name="Export PowerPoint").click()
        downloaded = tmp_path / "browser-native-visuals.pptx"
        download_event.value.save_as(downloaded)
        exported = Presentation(downloaded)
        assert all(any(shape.has_chart for shape in exported.slides[index].shapes) for index in (1, 2, 3))
        assert any(shape.name == "semantic.image.decision-flow" for shape in exported.slides[4].shapes)

        page.set_viewport_size({"width": 320, "height": 900})
        _assert_no_horizontal_overflow(page)
        assert page.get_by_label("Slide 1 chart type").is_visible()
        browser.close()


def test_accessibility_errors_reduced_motion_and_import_recovery(studio_url: str, tmp_path: Path):
    with sync_playwright() as playwright:
        browser = playwright.chromium.launch()
        page = browser.new_page(viewport={"width": 1280, "height": 900})
        page.goto(studio_url, wait_until="domcontentloaded")

        page.get_by_role("button", name="Build decision story").click()
        error = page.get_by_role("alert")
        assert error.is_visible()
        assert "decision" in error.inner_text().lower()
        assert page.locator(":focus").get_attribute("id") == "decision"

        unnamed_controls = page.locator("input, textarea, select, button").evaluate_all(
            """elements => elements
              .filter(element => element.type !== "hidden" && !element.hidden)
              .filter(element => {
                const labelled = element.getAttribute("aria-label")
                  || element.getAttribute("aria-labelledby")
                  || (element.id && document.querySelector(`label[for="${element.id}"]`))
                  || element.closest("label")
                  || element.textContent.trim();
                return !labelled;
              })
              .map(element => element.outerHTML.slice(0, 120))"""
        )
        assert unnamed_controls == []

        contrast = page.evaluate(
            """() => {
              const rgb = value => {
                const hex = value.replace("#", "");
                return [0, 2, 4].map(index => parseInt(hex.slice(index, index + 2), 16) / 255);
              };
              const luminance = value => rgb(value)
                .map(channel => channel <= .03928 ? channel / 12.92 : ((channel + .055) / 1.055) ** 2.4)
                .reduce((sum, channel, index) => sum + channel * [.2126, .7152, .0722][index], 0);
              const ratio = (first, second) => {
                const values = [luminance(first), luminance(second)].sort((a, b) => b - a);
                return (values[0] + .05) / (values[1] + .05);
              };
              return {
                body: ratio("#182321", "#f4f0e8"),
                action: ratio("#fffdf8", "#284a40"),
                muted: ratio("#5e6d67", "#fffdf8"),
              };
            }"""
        )
        assert min(contrast.values()) >= 4.5

        page.get_by_role("button", name="Try a sample brief").click()
        page.get_by_role("button", name="Build decision story").click()
        page.locator("#previewSection").wait_for(state="visible")
        invalid_outline = tmp_path / "invalid-outline.json"
        invalid_outline.write_text('{"title":"Missing slides"}', encoding="utf-8")
        page.locator("#importOutlineInput").set_input_files(invalid_outline)
        expect(page.locator("#saveStatus")).to_contain_text("Invalid outline")

        reduced_context = browser.new_context(reduced_motion="reduce")
        reduced_page = reduced_context.new_page()
        reduced_page.goto(studio_url, wait_until="domcontentloaded")
        reduced_state = reduced_page.evaluate(
            """() => ({
              matches: matchMedia("(prefers-reduced-motion: reduce)").matches,
              scrollBehavior: getComputedStyle(document.documentElement).scrollBehavior,
              transition: getComputedStyle(document.querySelector(".hero")).transitionDuration,
            })"""
        )
        assert reduced_state["matches"] is True
        assert reduced_state["scrollBehavior"] == "auto"
        transition_seconds = float(reduced_state["transition"].removesuffix("s"))
        assert transition_seconds <= 0.01
        reduced_context.close()
        browser.close()


def test_shared_layout_canvas_overflow_fixes_and_local_brand_kit(studio_url: str):
    with sync_playwright() as playwright:
        browser = playwright.chromium.launch()
        page = browser.new_page(viewport={"width": 1280, "height": 950})
        page.goto(studio_url, wait_until="domcontentloaded")
        page.locator('html[data-layout-ready="true"]').wait_for()
        assert "Shared layout v2" in page.locator("#layoutContractStatus").inner_text()

        page.get_by_role("button", name="Try a sample brief").click()
        page.get_by_role("button", name="Build decision story").click()
        page.locator("#previewSection").wait_for(state="visible")
        expect(page.locator("#deckPreview")).to_have_attribute("data-view", "canvas")

        first_slide = page.locator(".slide-preview").first
        size = first_slide.bounding_box()
        assert size is not None
        assert abs(size["width"] / size["height"] - 16 / 9) < 0.03
        assert page.locator(".preview-visual-label:visible").count() == 5

        page.set_viewport_size({"width": 320, "height": 900})
        expect(page.locator("#deckPreview")).to_have_attribute("data-view", "outline")
        _assert_no_horizontal_overflow(page)
        page.set_viewport_size({"width": 1280, "height": 950})
        expect(page.locator("#deckPreview")).to_have_attribute("data-view", "canvas")

        page.get_by_role("button", name="Zoom in").click()
        expect(page.locator("#zoomValue")).to_have_text("125%")
        page.get_by_role("button", name="Outline", exact=True).click()
        expect(page.locator("#deckPreview")).to_have_attribute("data-view", "outline")

        title = page.get_by_label("Slide 1 title")
        title.fill("This valid title is intentionally too long for the right layout now")
        title.press("Tab")
        expect(page.locator("#layoutPreflightTitle")).to_contain_text("layout fix", timeout=5000)
        assert page.locator('.slide-preview[data-overflow="true"]').count() >= 1
        page.get_by_role("button", name="Shorten to 60 characters").click()
        expect(page.locator("#layoutPreflightTitle")).to_have_text("Layout ready", timeout=5000)

        page.locator("#importBrandKitInput").set_input_files("themes/brand-kit.example.json")
        expect(page.locator("#saveStatus")).to_contain_text("Northstar brand kit")
        branded = page.locator(".slide-preview").first
        assert (
            branded.evaluate("element => getComputedStyle(element).backgroundColor") == "rgb(247, 249, 252)"
        )
        assert branded.evaluate("element => getComputedStyle(element).fontFamily").startswith("Aptos")
        page.get_by_role("button", name="Clear brand kit").click()

        page.set_viewport_size({"width": 320, "height": 900})
        expect(page.locator("#deckPreview")).to_have_attribute("data-view", "outline")
        _assert_no_horizontal_overflow(page)
        browser.close()


def test_complete_evidence_workflow_survives_edit_reorder_json_and_pptx(studio_url: str, tmp_path: Path):
    fixture = Path("examples/fixtures/evidence-edge-cases.json").resolve()
    with sync_playwright() as playwright:
        browser = playwright.chromium.launch()
        page = browser.new_page(viewport={"width": 1360, "height": 1000})
        page.goto(studio_url, wait_until="domcontentloaded")
        page.locator("#importOutlineInput").set_input_files(fixture)
        page.locator("#previewSection").wait_for(state="visible")

        expect(page.locator("#evidenceCoverageTitle")).to_contain_text("claims author checked")
        assert page.locator(".evidence-editor").count() == 3
        assert page.get_by_label("Source or evidence for slide 1").input_value() == "Étude d’usage — café"
        assert (
            page.get_by_label("Public URL for slide 1 source 1")
            .input_value()
            .startswith("https://example.org/")
        )
        assert page.get_by_label("Checked date for slide 1 source 1").input_value() == "2026-08-27"
        assert page.get_by_label("Review status for slide 1 source 1").input_value() == "author-checked"
        assert page.get_by_label("Slide summary for slide 1 source 1").is_checked()
        expect(page.get_by_role("button", name="Remove citations slide")).to_have_attribute(
            "aria-pressed", "true"
        )
        assert page.locator(".appendix-preview").count() == 1
        assert "Unreviewed résumé link" not in page.locator(".appendix-preview").inner_text()

        first_slide = page.locator(".content-preview").first
        first_slide.get_by_role("button", name="Duplicate slide").click()
        duplicated_label = page.get_by_label("Source or evidence for slide 2")
        duplicated_label.fill("Duplicated source preserved")
        duplicated_label.press("Tab")
        page.locator(".content-preview").nth(1).get_by_role("button", name="Move slide down").click()
        expect(page.get_by_label("Source or evidence for slide 3")).to_have_value(
            "Duplicated source preserved"
        )

        with page.expect_download() as json_download:
            page.get_by_role("button", name="Export JSON").click()
        story_path = tmp_path / "evidence.story.json"
        json_download.value.save_as(story_path)
        story = json.loads(story_path.read_text(encoding="utf-8"))
        moved_source = story["presentation"]["slides"][2]["sources"][0]
        assert moved_source["label"] == "Duplicated source preserved"
        assert moved_source["claim_ids"] == ["summary", "block-1"]
        assert moved_source["checked_date"] == "2026-08-27"

        with page.expect_download() as pptx_download:
            page.get_by_role("button", name="Export PowerPoint").click()
        pptx_path = tmp_path / "evidence-browser.pptx"
        pptx_download.value.save_as(pptx_path)
        exported = Presentation(pptx_path)
        appendix_text = "\n".join(shape.text for shape in exported.slides[-1].shapes if shape.has_text_frame)
        assert len(exported.slides) == 6
        assert "Duplicated source preserved" in appendix_text
        assert "Unreviewed résumé link" not in appendix_text

        page.get_by_role("button", name="Run Narrative Doctor").click()
        expect(page.locator("#doctorFindings")).to_contain_text("evidence.owner-missing")
        browser.close()


def test_markdown_story_and_local_source_excerpt_keep_claim_boundaries(studio_url: str, tmp_path: Path):
    story = Path("storyboard_studio/data/decision-brief.story.json").resolve()
    source_material = tmp_path / "onboarding-notes.txt"
    source_material.write_text(
        "Interview boundary one.\nThe pilot reduced handoff ambiguity for the reviewed cohort.\n"
        "Interview boundary three.\n",
        encoding="utf-8",
    )
    with sync_playwright() as playwright:
        browser = playwright.chromium.launch()
        page = browser.new_page(viewport={"width": 1360, "height": 1000})
        page.goto(studio_url, wait_until="domcontentloaded")
        page.locator("#importOutlineInput").set_input_files(story)
        page.locator("#previewSection").wait_for(state="visible")

        with page.expect_download() as markdown_download:
            page.get_by_role("button", name="Export Markdown").click()
        markdown_path = tmp_path / "review.story.md"
        markdown_download.value.save_as(markdown_path)
        assert "storyboard:content-block" in markdown_path.read_text(encoding="utf-8")
        page.locator("#importOutlineInput").set_input_files(markdown_path)
        expect(page.locator("#saveStatus")).to_contain_text("Markdown story imported locally")

        page.locator("#sourceMaterialInput").set_input_files(source_material)
        expect(page.locator("#sourceMaterialStatus")).to_contain_text("Loaded locally")
        excerpt = "The pilot reduced handoff ambiguity for the reviewed cohort."
        page.locator("#sourceMaterialText").evaluate(
            """(element, excerpt) => {
              const start = element.value.indexOf(excerpt);
              element.focus();
              element.setSelectionRange(start, start + excerpt.length);
              element.dispatchEvent(new Event('select', { bubbles: true }));
            }""",
            excerpt,
        )
        expect(page.locator("#sourceMaterialBoundary")).to_contain_text("lines 2–2")
        page.locator("#sourceMaterialSlide").select_option("0")
        page.locator("#sourceMaterialClaim").select_option("block-1")
        page.get_by_role("button", name="Map selected excerpt").click()

        expect(page.locator("#sourceMaterialStatus")).to_contain_text("slide 1 claim block-1")
        mapped = page.get_by_label("Label for slide 1 source 2")
        expect(mapped).to_have_value("onboarding-notes")
        expect(page.get_by_label("Local reference for slide 1 source 2")).to_have_value(
            "source-material/onboarding-notes.txt#L2-L2"
        )
        expect(page.get_by_label("Block claim 1 for slide 1 source 2")).to_be_checked()

        with page.expect_download() as mapped_markdown_download:
            page.get_by_role("button", name="Export Markdown").click()
        mapped_markdown = tmp_path / "mapped.story.md"
        mapped_markdown_download.value.save_as(mapped_markdown)
        assert excerpt in mapped_markdown.read_text(encoding="utf-8")
        assert "#L2-L2" in mapped_markdown.read_text(encoding="utf-8")
        browser.close()


def test_provider_is_disclosed_before_and_after_local_generation(studio_url: str):
    with sync_playwright() as playwright:
        browser = playwright.chromium.launch()
        page = browser.new_page(viewport={"width": 1280, "height": 900})
        external_requests: list[str] = []
        page.on(
            "request",
            lambda request: (
                external_requests.append(request.url) if not request.url.startswith(studio_url) else None
            ),
        )
        page.goto(studio_url, wait_until="domcontentloaded")

        expect(page.locator("#providerDisclosureTitle")).to_contain_text("deterministic-v1")
        expect(page.locator("#providerDisclosureTitle")).to_contain_text("offline")
        expect(page.locator("#providerDisclosureStatus")).to_contain_text("Fallback: not applicable")
        expect(page.locator("#providerDisclosurePolicy")).to_contain_text("evidence")
        assert page.locator('#providerSelect option[value="gemini"]').is_disabled()
        assert page.locator('#providerSelect option[value="openai-compatible"]').is_disabled()

        page.locator('input[name="workflow"][value="freeform"]').check()
        page.get_by_label("Presentation topic").fill("A local provider provenance review")
        page.get_by_role("button", name="Build my storyboard").click()
        page.locator("#previewSection").wait_for(state="visible")

        expect(page.locator("#providerRunTitle")).to_contain_text("deterministic-v1")
        expect(page.locator("#providerRunSummary")).to_contain_text("Selected local; used local")
        expect(page.locator("#providerRunSummary")).to_contain_text("Network: offline")
        expect(page.locator("#providerRunSummary")).to_contain_text("No fallback")
        expect(page.locator("#providerRunPolicy")).to_contain_text("sources")
        assert external_requests == []
        browser.close()
